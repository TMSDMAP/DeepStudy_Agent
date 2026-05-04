import html
import io
import json
import math
import os
import re
import shutil
import subprocess
import sys
import time
import ast
import hashlib
import base64
from datetime import datetime

import streamlit as st
from docx import Document
from docx.shared import Pt as DocxPt
from pptx.dml.color import RGBColor
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from workflow import app_workflow

try:
    import chromadb  # type: ignore
    HAS_CHROMA = True
except Exception:
    HAS_CHROMA = False


PPT_THEMES = {
    "teaching_blue": {
        "name": "课堂深蓝",
        "cover_bg": (12, 22, 44),
        "cover_deco": (16, 185, 129),
        "cover_glow": (31, 64, 104),
        "cover_subtitle": (186, 198, 217),
        "page_bg": (242, 246, 252),
        "header_bg": (17, 48, 89),
        "card_bg": (255, 255, 255),
        "card_border": (210, 220, 236),
        "title_text": (245, 250, 255),
        "body_text": (38, 53, 77),
        "code_text": (27, 39, 56),
    },
    "academic_gray": {
        "name": "学术灰白",
        "cover_bg": (31, 41, 55),
        "cover_deco": (20, 184, 166),
        "cover_glow": (70, 85, 108),
        "cover_subtitle": (204, 214, 224),
        "page_bg": (248, 250, 252),
        "header_bg": (51, 65, 85),
        "card_bg": (255, 255, 255),
        "card_border": (226, 232, 240),
        "title_text": (248, 250, 252),
        "body_text": (55, 65, 81),
        "code_text": (39, 50, 66),
    },
    "forest_class": {
        "name": "森系课堂",
        "cover_bg": (20, 45, 36),
        "cover_deco": (115, 187, 104),
        "cover_glow": (52, 92, 74),
        "cover_subtitle": (196, 216, 204),
        "page_bg": (241, 247, 243),
        "header_bg": (31, 74, 58),
        "card_bg": (255, 255, 255),
        "card_border": (204, 223, 211),
        "title_text": (245, 254, 248),
        "body_text": (45, 73, 60),
        "code_text": (34, 61, 49),
    },
}

def inject_app_styles():
    st.markdown(
        """
        <style>
        :root {
            --ink: #1f2937;
            --ink-soft: #4b5563;
            --panel: rgba(255, 255, 255, 0.92);
            --line: rgba(31, 41, 55, 0.14);
            --brand: #0f766e;
            --brand-soft: #14b8a6;
        }

        #MainMenu {visibility: hidden;}
        footer {visibility: hidden;}
        header {visibility: hidden;}
        .stDeployButton {display: none !important;}

        div[data-testid="stAppViewContainer"],
        .stApp {
            background:
                radial-gradient(38rem 24rem at 8% -12%, rgba(20, 184, 166, 0.18), transparent 70%),
                radial-gradient(34rem 28rem at 94% 8%, rgba(15, 118, 110, 0.15), transparent 72%),
                linear-gradient(180deg, #f8fafc 0%, #f3f4f6 100%);
        }

        .block-container {
            max-width: 1260px !important;
            padding-top: 2.1rem !important;
            padding-bottom: 3rem !important;
        }

        .hero-board {
            display: grid;
            grid-template-columns: 1.3fr 1fr;
            gap: 1rem;
            margin-bottom: 1.1rem;
        }

        .hero-main,
        .hero-side {
            border: 1px solid var(--line);
            border-radius: 22px;
            background: var(--panel);
            box-shadow: 0 14px 34px rgba(15, 23, 42, 0.08);
        }

        .hero-main { padding: 1.6rem 1.7rem; }
        .hero-side { padding: 1.1rem 1.2rem; }

        .hero-kicker {
            color: var(--brand);
            font-weight: 700;
            letter-spacing: .08em;
            text-transform: uppercase;
            font-size: .8rem;
            margin-bottom: .7rem;
        }

        .hero-title {
            margin: 0;
            color: #0f172a;
            font-size: clamp(1.95rem, 3.6vw, 3rem);
            line-height: 1.1;
            letter-spacing: -0.02em;
        }

        .hero-title span { color: var(--brand-soft); }

        .hero-desc {
            color: var(--ink-soft);
            margin-top: .9rem;
            margin-bottom: .95rem;
            line-height: 1.75;
        }

        .hero-tags {
            display: flex;
            flex-wrap: wrap;
            gap: .48rem;
        }

        .hero-tags span {
            border-radius: 999px;
            border: 1px solid rgba(15, 118, 110, 0.25);
            background: rgba(20, 184, 166, 0.08);
            color: #0f766e;
            padding: .22rem .72rem;
            font-size: .78rem;
            font-weight: 600;
        }

        .pipeline-title {
            margin: 0 0 .6rem 0;
            color: #64748b;
            font-size: .74rem;
            letter-spacing: .12em;
            text-transform: uppercase;
            font-weight: 700;
        }

        .pipeline-step {
            border: 1px solid rgba(100, 116, 139, 0.25);
            border-radius: 14px;
            background: rgba(255, 255, 255, 0.92);
            padding: .62rem .72rem;
            margin-bottom: .45rem;
            display: grid;
            grid-template-columns: 2rem 1fr;
            gap: .6rem;
        }

        .pipeline-step:last-child { margin-bottom: 0; }

        .pipeline-step b {
            width: 1.9rem;
            height: 1.9rem;
            border-radius: 10px;
            background: linear-gradient(145deg, #0f766e, #14b8a6);
            color: #fff;
            display: inline-flex;
            align-items: center;
            justify-content: center;
            font-size: .78rem;
        }

        .pipeline-step strong { color: #1f2937; }
        .pipeline-step p { margin: 0; color: #64748b; font-size: .82rem; }

        div[data-testid="stVerticalBlock"]:has(.prompt-shell-hook),
        div[data-testid="stVerticalBlock"]:has(.preset-shell-hook),
        div[data-testid="stVerticalBlock"]:has(.result-shell-hook),
        div[data-testid="stVerticalBlock"]:has(.code-audit-hook),
        div[data-testid="stVerticalBlock"]:has(.mode-shell-hook),
        div[data-testid="stVerticalBlock"]:has(.homework-shell-hook) {
            border-radius: 18px;
            border: 1px solid var(--line);
            background: rgba(255, 255, 255, 0.9);
            box-shadow: 0 10px 28px rgba(15, 23, 42, 0.08);
            padding: .95rem 1rem;
            margin-bottom: .9rem;
        }

        .preset-title,
        .export-title {
            color: #475569;
            font-size: .86rem;
            font-weight: 700;
            letter-spacing: .05em;
            text-transform: uppercase;
            margin-bottom: .45rem;
        }

        .result-head {
            display: flex;
            justify-content: space-between;
            align-items: center;
            padding: .2rem 0 .7rem;
            border-bottom: 1px solid rgba(100, 116, 139, 0.2);
            margin-bottom: .5rem;
        }

        .result-head-title {
            display: flex;
            align-items: center;
            gap: .52rem;
            color: #134e4a;
            font-size: .9rem;
            font-weight: 700;
        }

        .result-dot {
            width: 10px;
            height: 10px;
            border-radius: 999px;
            background: #14b8a6;
        }

        .result-meta {
            color: #64748b;
            font-size: .8rem;
            font-family: Consolas, monospace;
        }

        .code-audit-head {
            display: flex;
            justify-content: space-between;
            align-items: center;
            padding-bottom: .45rem;
            border-bottom: 1px dashed rgba(100, 116, 139, 0.25);
            margin-bottom: .65rem;
        }

        .code-audit-title { color: #134e4a; font-weight: 700; }
        .code-audit-meta { color: #64748b; font-size: .82rem; font-family: Consolas, monospace; }

        @media (max-width: 980px) {
            .hero-board { grid-template-columns: 1fr; }
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


def render_hero():
    st.markdown(
        """
        <section class="hero-board">
            <article class="hero-main">
                <div class="hero-kicker">DEEPLEARN COURSEWARE ENGINE</div>
                <h1 class="hero-title">深学讲义 <span>Agent</span></h1>
                <p class="hero-desc">
                    输入教学主题，多智能体自动完成知识检索、讲义生成、代码验证与课堂资料导出。
                    页面采用控制台 + 内容工作区结构，强调可读性与交付效率。
                </p>
                <div class="hero-tags">
                    <span>多 Agent 协同</span>
                    <span>代码执行质检</span>
                    <span>Markdown/Word/PDF/PPT/HTML</span>
                </div>
            </article>
            <aside class="hero-side">
                <h3 class="pipeline-title">Pipeline</h3>
                <div class="pipeline-step">
                    <b>01</b>
                    <div>
                        <strong>Subject Expert</strong>
                        <p>解析主题上下文，必要时调用搜索工具补充知识。</p>
                    </div>
                </div>
                <div class="pipeline-step">
                    <b>02</b>
                    <div>
                        <strong>Pedagogy Planner</strong>
                        <p>输出学情假设、课堂节奏与分层教学策略。</p>
                    </div>
                </div>
                <div class="pipeline-step">
                    <b>03</b>
                    <div>
                        <strong>Activity Designer</strong>
                        <p>设计互动任务、分组形式与形成性评价指标。</p>
                    </div>
                </div>
                <div class="pipeline-step">
                    <b>04</b>
                    <div>
                        <strong>Quiz Designer</strong>
                        <p>生成分层题库与答案解析，覆盖基础到进阶。</p>
                    </div>
                </div>
                <div class="pipeline-step">
                    <b>05</b>
                    <div>
                        <strong>Content Creator</strong>
                        <p>融合多 Agent 产物，生成课堂可用讲义。</p>
                    </div>
                </div>
                <div class="pipeline-step">
                    <b>06</b>
                    <div>
                        <strong>Evaluation Expert</strong>
                        <p>执行代码与格式审校，失败自动回写并重试。</p>
                    </div>
                </div>
            </aside>
        </section>
        """,
        unsafe_allow_html=True,
    )


def render_homework_hero():
    st.markdown(
        """
        <section class="hero-board">
            <article class="hero-main">
                <div class="hero-kicker">HOMEWORK INSIGHT ENGINE</div>
                <h1 class="hero-title">智能作答与批改 <span>Engine</span></h1>
                <p class="hero-desc">
                    上传题目或学生作答图片，自动完成题干数字化、变式生成、分步批改，并将结果沉淀到个人知识库。
                    布局与讲义引擎保持一致，便于教学流程切换。
                </p>
                <div class="hero-tags">
                    <span>多模态识别</span>
                    <span>变式生成</span>
                    <span>批改反馈</span>
                </div>
            </article>
            <aside class="hero-side">
                <h3 class="pipeline-title">Pipeline</h3>
                <div class="pipeline-step">
                    <b>01</b>
                    <div>
                        <strong>Image Intake</strong>
                        <p>上传或拍摄题目与作答图像。</p>
                    </div>
                </div>
                <div class="pipeline-step">
                    <b>02</b>
                    <div>
                        <strong>Vision Reasoning</strong>
                        <p>执行结构化提取、变式生成与逐步批改。</p>
                    </div>
                </div>
                <div class="pipeline-step">
                    <b>03</b>
                    <div>
                        <strong>Knowledge Sink</strong>
                        <p>将高质量结果沉淀到长期知识库。</p>
                    </div>
                </div>
            </aside>
        </section>
        """,
        unsafe_allow_html=True,
    )


def render_result_header(topic: str):
    generated_at = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    st.markdown(
        f"""
        <div class="result-head">
            <div class="result-head-title">
                <span class="result-dot"></span>
                教研大模型引擎输出内容
            </div>
            <div class="result-meta">{html.escape(topic)} · {generated_at}</div>
        </div>
        """,
        unsafe_allow_html=True,
    )


def strip_internal_tool_markup(text: str) -> str:
    if not text:
        return ""
    text = re.sub(r"</?[^>\n]*DSML[^>\n]*>", "", text, flags=re.I)
    cleaned_lines = []
    for line in text.splitlines():
        if "DSML" in line.upper():
            continue
        cleaned_lines.append(line)
    return "\n".join(cleaned_lines).strip()


INVISIBLE_CHARS = "\ufeff\u200b\u200c\u200d\u2060"


def _remove_invisible_chars(text: str) -> str:
    if not text:
        return ""
    cleaned = text
    for ch in INVISIBLE_CHARS:
        cleaned = cleaned.replace(ch, "")
    return cleaned


def _strip_quiz_inline_markup(text: str) -> str:
    if not text:
        return ""
    return text.replace("**", "").replace("__", "").replace("`", "").strip()


def normalize_quiz_markdown(text: str) -> str:
    lines = _remove_invisible_chars(text).splitlines()
    normalized: list[str] = []
    in_quiz_section = False
    in_answer_section = False
    answer_index = 1
    question_index = 1
    pending_answer_letter: str | None = None

    for raw in lines:
        line = _remove_invisible_chars(raw.rstrip())
        stripped = line.strip()
        match_text = _strip_quiz_inline_markup(stripped)

        if (stripped.startswith("##") or re.match(r"^[一二三四五六七八九十]+、", stripped)) and re.search(r"(课后选择题|练习题目)", stripped):
            in_quiz_section = True
            in_answer_section = False
            normalized.append("## 七、课后选择题（至少4题）")
            continue

        if (stripped.startswith("##") or re.match(r"^[一二三四五六七八九十]+、", stripped)) and re.search(r"答案", stripped):
            in_quiz_section = False
            in_answer_section = True
            answer_index = 1
            pending_answer_letter = None
            normalized.append("## 八、答案与解析")
            continue

        if (
            stripped.startswith("## ")
            and not re.search(r"(课后选择题|练习题目)", stripped)
            and not re.search(r"答案", stripped)
        ):
            in_quiz_section = False
            in_answer_section = False

        if in_quiz_section and re.fullmatch(r"\d+", match_text):
            # 清理偶发的孤立序号噪声（例如单独一行“2”）
            continue

        if in_quiz_section and re.match(r"^\d+\.\s+", match_text):
            q_num = re.match(r"^(\d+)\.\s+", match_text)
            if q_num:
                question_index = int(q_num.group(1)) + 1
            if all(re.search(rf"{opt}[\.．:：]\s*", match_text) for opt in ["A", "B", "C", "D"]):
                first_opt = re.search(r"[ABCD][\.．:：]\s*", match_text)
                stem = match_text[:first_opt.start()].strip() if first_opt else match_text
                normalized.append(stem)
                option_chunks = re.findall(
                    r"([ABCD])[\.．:：]\s*(.*?)(?=\s+[ABCD][\.．:：]\s*|$)",
                    match_text[first_opt.start():] if first_opt else "",
                )
                for opt, body in option_chunks:
                    normalized.append(f"{opt}. {body.strip()}")
                continue
            normalized.append(match_text)
            continue

        if in_quiz_section and all(re.search(rf"{opt}[\.．:：]\s*", match_text) for opt in ["A", "B", "C", "D"]):
            # 兼容“题干 + A/B/C/D 同行”但题干未编号
            first_opt = re.search(r"[ABCD][\.．:：]\s*", match_text)
            stem = match_text[:first_opt.start()].strip() if first_opt else match_text
            if stem:
                normalized.append(f"{question_index}. {stem}")
                question_index += 1
            option_chunks = re.findall(
                r"([ABCD])[\.．:：]\s*(.*?)(?=\s+[ABCD][\.．:：]\s*|$)",
                match_text[first_opt.start():] if first_opt else "",
            )
            for opt, body in option_chunks:
                normalized.append(f"{opt}. {body.strip()}")
            continue

        if in_quiz_section and re.match(r"^[\-*]\s*[ABCD][\.．:：]\s+", match_text):
            opt_line = re.sub(r"^[\-*]\s*", "", match_text)
            opt_match = re.match(r"^([ABCD])[\.．:：]\s*(.*)$", opt_line)
            if opt_match:
                normalized.append(f"{opt_match.group(1)}. {opt_match.group(2).strip()}")
            else:
                normalized.append(opt_line)
            continue

        if in_quiz_section and re.match(r"^[ABCD][\.．:：]\s+", match_text):
            opt_match = re.match(r"^([ABCD])[\.．:：]\s*(.*)$", match_text)
            if opt_match:
                normalized.append(f"{opt_match.group(1)}. {opt_match.group(2).strip()}")
            else:
                normalized.append(match_text)
            continue

        if in_answer_section:
            single_option = re.fullmatch(r"[A-Da-d]", match_text)
            numbered_option = re.fullmatch(r"(\d+)\.\s*([A-Da-d])", match_text)
            answer_with_explain = re.fullmatch(r"([A-Da-d])[\.．:：]\s*(.+)", match_text)
            numbered_with_explain = re.fullmatch(r"(\d+)\.\s*([A-Da-d])[\.．:：]\s*(.+)", match_text)
            correct_with_explain = re.fullmatch(
                r"(?:正确答案|答案)\s*[：:]\s*([A-Da-d])\s*(?:[，,；;。]?\s*)?(?:解析)\s*[：:]\s*(.+)",
                match_text,
            )
            answer_only = re.fullmatch(r"(?:正确答案|答案)\s*[：:]\s*([A-Da-d])\s*", match_text)
            explain_only = re.fullmatch(r"(?:解析)\s*[：:]\s*(.+)", match_text)
            if single_option:
                normalized.append(f"{answer_index}. {match_text.upper()}")
                answer_index += 1
                pending_answer_letter = None
                continue
            if correct_with_explain:
                letter = correct_with_explain.group(1).upper()
                explain = correct_with_explain.group(2).strip()
                normalized.append(f"{answer_index}. {letter}：{explain}")
                answer_index += 1
                pending_answer_letter = None
                continue
            if numbered_option:
                normalized.append(f"{numbered_option.group(1)}. {numbered_option.group(2).upper()}")
                answer_index = int(numbered_option.group(1)) + 1
                pending_answer_letter = None
                continue
            if answer_with_explain:
                letter = answer_with_explain.group(1).upper()
                explain = answer_with_explain.group(2).strip()
                normalized.append(f"{answer_index}. {letter}：{explain}")
                answer_index += 1
                pending_answer_letter = None
                continue
            if numbered_with_explain:
                num = int(numbered_with_explain.group(1))
                letter = numbered_with_explain.group(2).upper()
                explain = numbered_with_explain.group(3).strip()
                normalized.append(f"{num}. {letter}：{explain}")
                answer_index = num + 1
                pending_answer_letter = None
                continue
            if answer_only:
                pending_answer_letter = answer_only.group(1).upper()
                continue
            if explain_only and pending_answer_letter:
                normalized.append(f"{answer_index}. {pending_answer_letter}：{explain_only.group(1).strip()}")
                answer_index += 1
                pending_answer_letter = None
                continue

        normalized.append(match_text if in_answer_section else line)

    return "\n".join(normalized).strip()


def _enforce_min_quiz_items(text: str, topic: str, min_questions: int = 4) -> str:
    quiz_block_match = re.search(
        r"(?ms)^##\s*七、课后选择题（至少4题）\s*(.*?)(?=^##\s*八、答案与解析|\Z)",
        text,
    )
    answer_block_match = re.search(
        r"(?ms)^##\s*八、答案与解析\s*(.*?)(?=^##\s|\Z)",
        text,
    )
    if not quiz_block_match or not answer_block_match:
        return text

    quiz_body = _remove_invisible_chars(quiz_block_match.group(1).strip())
    existing_questions = len(re.findall(r"(?m)^\d+\.\s+", quiz_body))
    if existing_questions == 0:
        existing_questions = len(
            re.findall(
                r"(?m)^.*A[\.．:：].*B[\.．:：].*C[\.．:：].*D[\.．:：].*$",
                quiz_body,
            )
        )
    needed_questions = max(min_questions, existing_questions)

    if existing_questions < min_questions:
        supplements = []
        for idx in range(existing_questions + 1, min_questions + 1):
            supplements.append(_fallback_quiz_item(topic, idx))
        quiz_body = (quiz_body + "\n\n" + "\n\n".join(supplements)).strip() if quiz_body else "\n\n".join(supplements)
        text = text[:quiz_block_match.start(1)] + quiz_body + text[quiz_block_match.end(1):]

    # 重新定位答案区块，确保索引正确
    answer_block_match = re.search(
        r"(?ms)^##\s*八、答案与解析\s*(.*?)(?=^##\s|\Z)",
        text,
    )
    if not answer_block_match:
        return text

    answer_body = _remove_invisible_chars(answer_block_match.group(1).strip())
    answer_lines = [_remove_invisible_chars(ln.strip()) for ln in answer_body.splitlines() if ln.strip()]
    answer_lines = [ln for ln in answer_lines if "待补充答案" not in ln]
    answer_body = "\n".join(answer_lines).strip()
    existing_answers = len(re.findall(r"(?m)^\d+\.\s+", answer_body))

    if existing_answers == 0:
        letter_answers = re.findall(r"(?m)^([A-Da-d])[\.．:：]\s*(.+)$", answer_body)
        if letter_answers:
            rebuilt_answers = []
            for idx, (letter, explain) in enumerate(letter_answers, start=1):
                rebuilt_answers.append(f"{idx}. {letter.upper()}：{explain.strip()}")
            answer_body = "\n".join(rebuilt_answers)
            existing_answers = len(rebuilt_answers)

    if existing_answers == 0:
        compact_answers = re.findall(
            r"(?m)(?:正确答案|答案)\s*[：:]\s*([A-Da-d])(?:\s*(?:[，,；;。]?\s*)?(?:解析)\s*[：:]\s*(.+))?",
            answer_body,
        )
        if compact_answers:
            rebuilt_answers = []
            for idx, (letter, explain) in enumerate(compact_answers, start=1):
                parsed_letter = (letter or "").upper()
                parsed_explain = (explain or "").strip()
                if parsed_explain:
                    rebuilt_answers.append(f"{idx}. {parsed_letter}：{parsed_explain}")
                else:
                    rebuilt_answers.append(f"{idx}. {parsed_letter}")
            answer_body = "\n".join(rebuilt_answers)
            existing_answers = len(rebuilt_answers)

    if existing_answers > needed_questions:
        numbered_lines = [ln for ln in answer_body.splitlines() if re.match(r"^\d+\.\s+", ln.strip())]
        answer_body = "\n".join(numbered_lines[:needed_questions]).strip()
        existing_answers = len(re.findall(r"(?m)^\d+\.\s+", answer_body))

    if existing_answers < needed_questions:
        extra_answers = [_fallback_answer_line(idx) for idx in range(existing_answers + 1, needed_questions + 1)]
        answer_body = (answer_body + "\n" + "\n".join(extra_answers)).strip() if answer_body else "\n".join(extra_answers)
        text = text[:answer_block_match.start(1)] + answer_body + text[answer_block_match.end(1):]
    else:
        text = text[:answer_block_match.start(1)] + answer_body + text[answer_block_match.end(1):]

    return text


def _replace_markdown_section(text: str, heading_regex: str, canonical_heading: str, body: str) -> str:
    section_pattern = rf"(?ms)^##\s*.*?(?:{heading_regex}).*?(?=^##\s|\Z)"
    replacement = f"{canonical_heading}\n{body.strip()}\n"
    if re.search(section_pattern, text):
        return re.sub(section_pattern, replacement, text, count=1).strip()
    return (text.rstrip() + "\n\n" + replacement).strip()


def _normalize_major_section_heading_markers(text: str) -> str:
    """Convert bare Chinese numbered headings to markdown H2 (e.g. '七、课后...' -> '## 七、课后...')."""
    lines = text.splitlines()
    normalized: list[str] = []
    for raw in lines:
        line = _remove_invisible_chars(raw)
        stripped = line.strip()
        if stripped.startswith("#"):
            normalized.append(line)
            continue
        if re.match(r"^[一二三四五六七八九十]+、", stripped):
            normalized.append(f"## {stripped}")
            continue
        normalized.append(line)
    return "\n".join(normalized)


def _looks_like_python_code(line: str) -> bool:
    stripped = line.strip()
    if not stripped:
        return False
    patterns = [
        r"^(import|from)\s+\w+",
        r"^def\s+\w+\(",
        r"^class\s+\w+",
        r"^(if|elif|else|for|while|try|except|with|return|yield|pass|break|continue)\b",
        r"^@\w+",
        r"^\w[\w\[\]\.]*\s*=\s*.+",
        r"^(print|len|range|list|dict|set|tuple|sum|min|max)\s*\(",
        r"^#",
    ]
    return any(re.match(p, stripped) for p in patterns)


def _auto_wrap_code_section(text: str) -> str:
    match = re.search(r"(?ms)^##\s*.*?(示例代码|代码演示).*?\n(.*?)(?=^##\s|\Z)", text)
    if not match:
        return text
    body = (match.group(2) or "").strip()
    if not body or "```" in body:
        return text

    lines = [ln.rstrip() for ln in body.splitlines()]
    narrative: list[str] = []
    code: list[str] = []
    code_started = False

    for ln in lines:
        stripped = ln.strip()
        if not code_started and _looks_like_python_code(stripped):
            code_started = True
        if code_started:
            code.append(ln)
        else:
            narrative.append(ln)

    code_nonempty = [ln for ln in code if ln.strip()]
    if len(code_nonempty) < 2:
        return text

    narrative_text = "\n".join(narrative).strip()
    code_text = "\n".join(code).rstrip()
    new_body = (narrative_text + "\n\n" if narrative_text else "") + f"```python\n{code_text}\n```"
    return text[:match.start(2)] + new_body + text[match.end(2):]


def _fallback_quiz_item(topic: str, idx: int) -> str:
    topic_name = topic or "本课主题"
    stems = {
        1: f"围绕“{topic_name}”的核心概念，下列说法最准确的是哪一项？",
        2: f"在“{topic_name}”的课堂实践环节中，下列哪项最符合教学目标对齐原则？",
        3: f"如果要检验学生是否真正理解“{topic_name}”，下列哪种任务设计更有效？",
        4: f"针对“{topic_name}”的常见误区，下列哪项纠偏策略更合理？",
    }
    opts = {
        1: [
            "A. 只记住术语定义，不需要理解应用场景",
            "B. 能解释概念内涵，并结合案例说明其作用边界",
            "C. 只要会背公式就说明掌握了概念",
            "D. 概念学习与课堂活动没有关系",
        ],
        2: [
            "A. 先讲结论再直接布置作业，不做过程引导",
            "B. 只展示标准答案，不安排学生表达思路",
            "C. 目标、活动、评价三者一致，并有可观察表现",
            "D. 课堂目标与评价标准分离，靠课后补救",
        ],
        3: [
            "A. 让学生复述教材原文即可",
            "B. 让学生完成迁移任务并解释决策依据",
            "C. 只做选择题且不要求解释",
            "D. 只看作业完成速度",
        ],
        4: [
            "A. 忽略错误，避免影响课堂节奏",
            "B. 只指出对错，不分析原因",
            "C. 让学生自行查资料，不给支架",
            "D. 定位错误来源并用对比样例完成再建构",
        ],
    }
    stem = stems.get(idx, f"关于“{topic_name}”的课堂理解，下列哪项最恰当？")
    options = opts.get(
        idx,
        [
            "A. 说法A",
            "B. 说法B",
            "C. 说法C",
            "D. 说法D",
        ],
    )
    return f"{idx}. {stem}\n" + "\n".join(options)


def _fallback_answer_line(idx: int) -> str:
    defaults = {
        1: ("B", "该选项能同时体现概念含义、适用场景和分析边界，比单纯记忆术语更符合课堂理解要求。"),
        2: ("C", "教学活动应与学习目标和评价证据保持一致，学生的课堂产出需要能够证明目标是否达成。"),
        3: ("B", "迁移任务要求学生解释依据，能检查其是否真正理解知识并能在新情境中应用。"),
        4: ("D", "有效纠偏需要先定位错误来源，再通过对比样例和再练习帮助学生重建理解。"),
    }
    letter, explanation = defaults.get(idx, ("B", "该选项最符合本课的核心理解要求。"))
    return f"{idx}. {letter}：{explanation}"


def _replace_placeholder_quiz_if_needed(text: str, topic: str, min_questions: int = 4) -> str:
    placeholder_hit = re.search(
        r"(请补充一道选择题|请补充第二道|^\d+\.\s*第三道$|^\d+\.\s*第四道$|A\.\s*选项A|^\d+\.\s*[A-D]：解析\s*$)",
        text,
        flags=re.M,
    )
    if not placeholder_hit:
        return text

    quiz_body = "\n\n".join(_fallback_quiz_item(topic, idx) for idx in range(1, min_questions + 1))
    answer_body = "\n".join(_fallback_answer_line(idx) for idx in range(1, min_questions + 1))
    text = _replace_markdown_section(text, r"(课后选择题|练习题目)", "## 七、课后选择题（至少4题）", quiz_body)
    text = _replace_markdown_section(text, r"答案", "## 八、答案与解析", answer_body)
    return text


def _remove_markdown_sections_by_keyword(text: str, keyword_regex: str) -> str:
    return re.sub(
        rf"(?ms)^##\s*.*?(?:{keyword_regex}).*?(?=^##\s|\Z)",
        "",
        text,
    ).strip()


def _extract_topic_scope_hint(topic: str) -> str:
    if not topic:
        return ""
    text = topic.strip()
    if re.search(r"(最后一回|末回|终回|第一百二十回|第120回)", text):
        return "最后一回（第一百二十回）"
    match = re.search(r"第\s*([一二三四五六七八九十百零两\d]+)\s*回", text)
    if match:
        return f"第{match.group(1)}回"
    return ""


def _enforce_topic_scope_guard(text: str, topic: str) -> str:
    scope = _extract_topic_scope_hint(topic)
    if not scope:
        return text

    if scope == "最后一回（第一百二十回）":
        first_count = len(re.findall(r"第一回", text))
        last_count = len(re.findall(r"(最后一回|第一百二十回|第120回|一百二十回)", text))
        if first_count == 0 or last_count >= first_count:
            return text
        # 当主题明确要求“最后一回”且正文完全未出现时，纠正常见的“第一回”漂移。
        text = re.sub(r"(?m)^(#+\s*.*)第一回", r"\1最后一回（第一百二十回）", text)
        text = re.sub(r"第一回", "最后一回（第一百二十回）", text, count=min(5, first_count))
        return text

    if scope not in text:
        text = re.sub(
            r"(?m)^(#+\s*.*)第[一二三四五六七八九十百零两\d]+回",
            lambda m: re.sub(r"第[一二三四五六七八九十百零两\d]+回", scope, m.group(0), count=1),
            text,
            count=1,
        )
    return text


def _expand_html_break_tags(text: str) -> str:
    if not text:
        return ""
    text = html.unescape(text)
    return re.sub(r"(?i)<br\s*/?>", "\n", text)


def ensure_markdown_template(md_text: str, topic: str, include_code: bool = True, include_quiz: bool = True) -> str:
    topic_title = (topic or "教学主题").strip() or "教学主题"
    text = _remove_invisible_chars(strip_internal_tool_markup(md_text)).strip()
    text = _expand_html_break_tags(text)
    text = _normalize_major_section_heading_markers(text)
    if not text:
        text = f"# {topic_title} 讲义\n"
    canonical_h1 = f"# {topic_title} 讲义"
    if re.search(r"(?m)^#\s+", text):
        text = re.sub(r"(?m)^#\s+.*$", canonical_h1, text, count=1)
    else:
        text = f"{canonical_h1}\n\n{text}"

    if not include_code:
        text = _remove_markdown_sections_by_keyword(text, r"(示例代码|代码演示)")
        text = re.sub(r"(?ms)```python.*?```", "", text)
    if not include_quiz:
        text = _remove_markdown_sections_by_keyword(text, r"(课后选择题|选择题|练习题目|答案)")

    required_sections = [
        (r"##\s*.*?学习目标", "## 一、学习目标", "- 目标1\n- 目标2"),
        (r"##\s*.*?(先修知识|学情)", "## 二、先修知识与学情假设", "- 先修知识1\n- 学情假设1"),
        (r"##\s*.*?(课堂流程|时间分配)", "## 三、课堂流程（含时间分配）", "- 导入（10分钟）\n- 讲授（25分钟）\n- 练习（15分钟）"),
        (r"##\s*.*?核心知识", "## 四、核心知识点", "### 4.1 概念\n### 4.2 原理\n### 4.3 常见误区"),
        (r"##\s*.*?互动活动", "## 五、互动活动设计", "- 活动1\n- 活动2"),
    ]
    if include_code:
        required_sections.append(
            (
                r"##\s*.*?(示例代码|代码演示)",
                "## 六、示例代码",
                "- 请补充与主题相关的可运行示例代码（建议包含输入、处理、输出）。",
            )
        )
    if include_quiz:
        required_sections.extend(
            [
                (
                    r"##\s*.*?(选择题|练习题目)",
                    "## 七、课后选择题（至少4题）",
                    "1. 请补充一道选择题\nA. 选项A\nB. 选项B\nC. 选项C\nD. 选项D\n2. 请补充第二道\nA. A\nB. B\nC. C\nD. D\n3. 第三道\nA. A\nB. B\nC. C\nD. D\n4. 第四道\nA. A\nB. B\nC. C\nD. D",
                ),
                (r"##\s*.*?答案", "## 八、答案与解析", "1. A：解析\n2. B：解析\n3. C：解析\n4. D：解析"),
            ]
        )

    for pattern, default_heading, default_body in required_sections:
        if not re.search(pattern, text):
            text += f"\n\n{default_heading}\n{default_body}\n"

    if include_code:
        text = _auto_wrap_code_section(text)

    if include_code and not extract_python_code_blocks(text):
        fallback_code = (
            "```python\n"
            f"lesson_topic = {topic_title!r}\n"
            "print(f\"本节主题：{lesson_topic}\")\n"
            "for idx, item in enumerate([\"核心概念\", \"关键问题\", \"课堂任务\"], start=1):\n"
            "    print(f\"{idx}. {item}\")\n"
            "```"
        )
        code_section_match = re.search(r"(?ms)^##\s*.*?(示例代码|代码演示).*?(?=^##\s|\Z)", text)
        if code_section_match:
            section_text = code_section_match.group(0).rstrip()
            replacement = f"{section_text}\n\n{fallback_code}\n"
            text = text[:code_section_match.start()] + replacement + text[code_section_match.end():]
        else:
            text += f"\n\n## 六、示例代码\n{fallback_code}\n"

    if include_quiz:
        text = normalize_quiz_markdown(text)
        text = _replace_placeholder_quiz_if_needed(text, topic, min_questions=4)
        text = _enforce_min_quiz_items(text, topic=topic, min_questions=4)
    text = _enforce_topic_scope_guard(text, topic_title)
    return text.strip()


def _lesson_output_quality_issues(md_text: str, include_code: bool = True, include_quiz: bool = True) -> list[str]:
    text = (md_text or "").strip()
    if include_quiz:
        text = normalize_quiz_markdown(text)
    compact = re.sub(r"\s+", "", text)
    issues: list[str] = []

    if len(compact) < 700:
        issues.append("讲义正文过短")

    required = [
        ("学习目标", r"##\s*.*学习目标"),
        ("先修知识与学情", r"##\s*.*(先修知识|学情)"),
        ("课堂流程", r"##\s*.*(课堂流程|时间分配)"),
        ("核心知识点", r"##\s*.*核心知识"),
        ("互动活动设计", r"##\s*.*互动活动"),
    ]
    missing = [name for name, pattern in required if not re.search(pattern, text)]
    if missing:
        issues.append("缺少主体章节：" + "、".join(missing))

    if re.search(r"(?i)(待补充答案|请补充|TODO|# TODO)", text):
        issues.append("仍包含占位内容")

    if include_code and not extract_python_code_blocks(text):
        issues.append("缺少可运行 python 代码块")

    if include_quiz:
        quiz_match = re.search(
            r"(?ms)^##\s*七、课后选择题（至少4题）\s*(.*?)(?=^##\s*八、答案与解析|\Z)",
            text,
        )
        answer_match = re.search(
            r"(?ms)^##\s*八、答案与解析\s*(.*?)(?=^##\s|\Z)",
            text,
        )
        quiz_body = quiz_match.group(1).strip() if quiz_match else ""
        answer_text = answer_match.group(1).strip() if answer_match else ""
        quiz_count = len(re.findall(r"(?m)^\d+\.\s+", quiz_body))
        if quiz_count == 0:
            quiz_count = len(
                re.findall(
                    r"(?m)^.*A[\.．:：].*B[\.．:：].*C[\.．:：].*D[\.．:：].*$",
                    quiz_body,
                )
            )
        answer_count = len(
            re.findall(
                r"(?m)^\s*(?:\d+\.\s*)?(?:[A-D][：:]\s*|(?:正确答案|答案)\s*[：:]\s*[A-D])",
                answer_text,
            )
        )
        if quiz_count < 4:
            issues.append("选择题不足 4 道")
        if answer_count < 4:
            issues.append("答案解析不足 4 条")

    return issues


def _collect_runtime_review_debug(md_text: str, include_quiz: bool = True) -> dict:
    text = _remove_invisible_chars(strip_internal_tool_markup(md_text or "")).strip()
    if include_quiz:
        text = normalize_quiz_markdown(text)

    quiz_match = re.search(
        r"(?ms)^##\s*七、课后选择题（至少4题）\s*(.*?)(?=^##\s*八、答案与解析|\Z)",
        text,
    )
    answer_match = re.search(
        r"(?ms)^##\s*八、答案与解析\s*(.*?)(?=^##\s|\Z)",
        text,
    )
    quiz_text = quiz_match.group(1).strip() if quiz_match else ""
    answer_text = answer_match.group(1).strip() if answer_match else ""

    question_count = len(re.findall(r"(?m)^\d+\.\s+", quiz_text))
    if question_count == 0:
        question_count = len(
            re.findall(
                r"(?m)^.*A[\.．:：].*B[\.．:：].*C[\.．:：].*D[\.．:：].*$",
                quiz_text,
            )
        )

    answer_count = len(
        re.findall(
            r"(?m)^\s*(?:\d+\.\s*)?(?:[A-D][：:]\s*|(?:正确答案|答案)\s*[：:]\s*[A-D])",
            answer_text,
        )
    )

    quiz_preview = "\n".join(quiz_text.splitlines()[:12]).strip()
    answer_preview = "\n".join(answer_text.splitlines()[:12]).strip()
    return {
        "question_count": question_count,
        "answer_count": answer_count,
        "quiz_preview": quiz_preview,
        "answer_preview": answer_preview,
    }


def _brief_agent_block(text: str, fallback: str, max_chars: int = 2800) -> str:
    cleaned = _remove_invisible_chars(strip_internal_tool_markup(text or "")).strip()
    if len(cleaned) < 40:
        return fallback.strip()
    return cleaned[:max_chars].strip()


def _fallback_demo_code_block(topic: str) -> str:
    safe_topic = (topic or "本课主题").strip() or "本课主题"
    return f"""```python
lesson_topic = {safe_topic!r}
key_points = [\"核心概念\", \"关键问题\", \"课堂任务\"]
print(f\"本节主题：{{lesson_topic}}\")
for idx, point in enumerate(key_points, start=1):
    print(f\"{{idx}}. {{point}}\")
```"""


def _build_recovery_lesson_markdown(topic: str, state: dict, include_code: bool, include_quiz: bool) -> str:
    outline = _brief_agent_block(
        state.get("outline", ""),
        f"本课围绕“{topic}”展开，重点覆盖核心概念、关键问题、常见误区与课堂应用场景。",
    )
    pedagogy = _brief_agent_block(
        state.get("pedagogy_plan", ""),
        "学生已有基础知识但理解深度不均。课堂需要通过分层提问、可视化示例与过程性评价帮助学生建立结构化认知。",
    )
    activity = _brief_agent_block(
        state.get("activity_plan", ""),
        "活动建议：先进行小组任务拆解，再做全班汇报与同伴互评，最后结合教师点评完成知识回收与纠偏。",
    )

    parts = [
        f"# {topic} 讲义",
        "## 一、学习目标",
        f"本课目标是让学生形成对“{topic}”的清晰理解，并能够在课堂任务中准确表达、分析和迁移应用。",
        "- 能用自己的语言解释核心概念与关键术语。",
        "- 能结合课堂案例说明概念的适用边界。",
        "- 能识别常见误区并给出修正理由。",
        "- 能在小组任务中完成证据化表达。",
        "- 能把所学迁移到新情境进行判断。",
        "## 二、先修知识与学情假设",
        pedagogy,
        "## 三、课堂流程（含时间分配）",
        "建议按 45 分钟组织：导入 6 分钟、核心讲解 16 分钟、互动活动 14 分钟、练习反馈 7 分钟、总结与布置 2 分钟。",
        "- 导入：问题情境唤醒已有经验，明确本课任务。",
        "- 讲解：结合示例建立概念框架与关键联系。",
        "- 活动：分组任务 + 汇报互评，沉淀学习证据。",
        "- 反馈：针对误区进行对比纠偏与再练。",
        "- 总结：回收目标并给出课后迁移任务。",
        "## 四、核心知识点",
        outline,
        "## 五、互动活动设计",
        activity,
    ]

    if include_code:
        parts.extend(
            [
                "## 六、示例代码",
                "以下代码用于课堂演示如何用结构化输出梳理知识点，便于学生观察输入、处理与输出过程。",
                _fallback_demo_code_block(topic),
            ]
        )

    if include_quiz:
        parts.extend(
            [
                "## 七、课后选择题（至少4题）",
                "\n\n".join(_fallback_quiz_item(topic, idx) for idx in range(1, 5)),
                "## 八、答案与解析",
                "\n".join(_fallback_answer_line(idx) for idx in range(1, 5)),
            ]
        )

    return "\n\n".join(parts).strip()

def _expand_inline_quiz_options(text: str) -> str:
    if not text:
        return ""
    lines = text.splitlines()
    expanded: list[str] = []
    in_code_block = False

    for raw in lines:
        line = raw.rstrip()
        stripped = line.strip()

        if stripped.startswith("```"):
            in_code_block = not in_code_block
            expanded.append(line)
            continue

        if in_code_block or not stripped:
            expanded.append(line)
            continue

        option_marks = list(re.finditer(r"([ABCD])[\.．:：]\s*", stripped))
        option_letters = {m.group(1) for m in option_marks}
        if len(option_marks) >= 4 and option_letters.issuperset({"A", "B", "C", "D"}):
            first_opt = option_marks[0]
            stem = stripped[:first_opt.start()].strip()
            option_chunks = re.findall(
                r"([ABCD])[\.．:：]\s*(.*?)(?=\s+[ABCD][\.．:：]\s*|$)",
                stripped[first_opt.start():],
            )
            if stem:
                expanded.append(stem)
            for opt, body in option_chunks:
                expanded.append(f"{opt}. {body.strip()}")
            continue

        expanded.append(line)

    return "\n".join(expanded)


def _is_schedule_table_header(line: str) -> bool:
    compact = re.sub(r"\s+", "", line or "")
    return all(token in compact for token in ("阶段", "时间", "教学活动与步骤"))


def _is_generic_tabular_header(line: str) -> bool:
    stripped = (line or "").strip()
    if not stripped:
        return False
    if re.match(r"^\d+[\.、]", stripped):
        return False
    if "\t" in stripped:
        cells = [cell.strip() for cell in re.split(r"\t+", stripped) if cell.strip()]
        if len(cells) >= 3 and len(cells[0]) <= 8 and len(cells[1]) <= 8:
            return True
    if "|" in stripped and stripped.count("|") >= 2:
        raw_cells = [cell.strip() for cell in stripped.strip("|").split("|")]
        cells = [cell for cell in raw_cells if cell]
        if len(cells) >= 3 and len(cells[0]) <= 8 and len(cells[1]) <= 8:
            return True
    return False


def _parse_schedule_table_row(line: str) -> tuple[str, str, str] | None:
    stripped = (line or "").strip()
    if not stripped:
        return None

    if re.fullmatch(r"[\|\-:\s]+", stripped):
        return None

    if "\t" in stripped:
        cells = [cell.strip() for cell in re.split(r"\t+", stripped) if cell.strip()]
    elif "|" in stripped and stripped.count("|") >= 2:
        raw_cells = [cell.strip() for cell in stripped.strip("|").split("|")]
        cells = [cell for cell in raw_cells if cell]
    else:
        cells = [cell.strip() for cell in re.split(r"\s{2,}", stripped) if cell.strip()]

    if len(cells) >= 3 and re.match(r"^\d+[\.、]", cells[0]) and re.search(r"\d+", cells[1]):
        return cells[0], cells[1], " ".join(cells[2:])

    match = re.match(r"^(\d+[\.、]\s*.+?)\s+(\d+\s*\S*)\s+(.+)$", stripped)
    if match:
        stage = match.group(1).strip()
        duration = match.group(2).strip()
        body = match.group(3).strip()
        return stage, duration, body
    return None


def _normalize_schedule_table_block(text: str) -> str:
    if not text:
        return ""
    lines = text.splitlines()
    normalized: list[str] = []
    in_code_block = False
    i = 0

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if stripped.startswith("```"):
            in_code_block = not in_code_block
            normalized.append(line)
            i += 1
            continue

        if in_code_block:
            normalized.append(line)
            i += 1
            continue

        header_hit = _is_schedule_table_header(stripped)
        seed_row = _parse_schedule_table_row(stripped) if not header_hit else None

        if not header_hit and not seed_row:
            normalized.append(line)
            i += 1
            continue

        j = i + 1 if header_hit else i + 1
        parsed_rows: list[tuple[str, str, str]] = []
        if seed_row:
            parsed_rows.append(seed_row)

        while j < len(lines):
            candidate = lines[j].strip()
            if not candidate:
                j += 1
                if parsed_rows:
                    break
                continue
            if re.fullmatch(r"[\|\-:\s]+", candidate):
                j += 1
                continue
            parsed = _parse_schedule_table_row(candidate)
            if not parsed:
                break
            parsed_rows.append(parsed)
            j += 1

        if not parsed_rows:
            normalized.append(line)
            i += 1
            continue

        if seed_row and normalized and _is_generic_tabular_header(normalized[-1]):
            normalized.pop()

        for stage, duration, body in parsed_rows:
            normalized.append(f"### {stage}（{duration}）")
            body_text = _expand_html_break_tags(body)
            segments = [seg.strip() for seg in body_text.splitlines() if seg.strip()]
            if not segments:
                normalized.append("- （内容待补充）")
            else:
                for seg in segments:
                    normalized.append(f"- {seg}")
            normalized.append("")

        i = j

    return "\n".join(normalized)


def sanitize_markdown_for_ppt(md_text: str) -> str:
    """Remove layout-hallucination markers before PPT rendering."""
    if not md_text:
        return ""
    md_text = _remove_invisible_chars(md_text)
    md_text = _normalize_schedule_table_block(md_text)
    md_text = _expand_html_break_tags(md_text)
    md_text = _expand_inline_quiz_options(md_text)

    bad_words = {
        "dual column",
        "本页摘要",
        "结构化模板",
        "多版式分页",
        "课堂可直接投屏",
        "关键提炼",
        "结构化讲解",
        "核心要点",
        "补充要点",
        "主题导读",
        "重点解析",
        "课堂演示",
        "本节要点",
        "教学内容",
        "主要内容",
        "延展内容",
    }

    punctuation_pattern = re.compile(r"[，。！？：；,.!?:;（）()【】\[\]《》<>]")

    lines = md_text.splitlines()
    cleaned_lines = []
    for line in lines:
        stripped = line.strip()
        stripped_key = re.sub(r"^[#\-\*\s]+|[#\-\*\s]+$", "", stripped).lower()

        if stripped_key in bad_words or re.fullmatch(r"要点\s*[A-Z]?", stripped_key, flags=re.IGNORECASE):
            continue

        if re.match(r"^.*代码示例.*段.*$", stripped):
            continue

        if not stripped.startswith("#") and not stripped.startswith("```"):
            clean_text = re.sub(r"[\-\*\•\s]", "", stripped)
            if (
                2 <= len(clean_text) <= 5
                and re.match(r"^[\u4e00-\u9fa5]+$", clean_text)
                and not punctuation_pattern.search(stripped)
            ):
                continue

        cleaned_lines.append(line)

    return "\n".join(cleaned_lines).strip()


def _text_quality_score(text: str) -> int:
    if not text:
        return -10_000
    cjk_count = len(re.findall(r"[\u4e00-\u9fff]", text))
    mojibake_count = len(re.findall(r"[ÃÂâãæçèéêäåœ€]", text))
    replacement_count = text.count("\ufffd")
    return cjk_count * 4 - mojibake_count * 3 - replacement_count * 80


def _repair_mojibake_text(text: str) -> str:
    if not text:
        return ""
    original_score = _text_quality_score(text)
    best = text
    best_score = original_score
    for enc in ("cp1252", "latin1"):
        try:
            repaired = text.encode(enc, errors="strict").decode("utf-8", errors="strict")
        except Exception:
            continue
        repaired_score = _text_quality_score(repaired)
        if repaired_score > best_score + 20:
            best = repaired
            best_score = repaired_score
    return best


def _decode_text_bytes(raw: bytes) -> str:
    candidates: list[str] = []
    for enc in ("utf-8-sig", "utf-8", "gb18030", "gbk", "big5"):
        try:
            decoded = raw.decode(enc)
        except Exception:
            continue
        candidates.append(_repair_mojibake_text(decoded))

    if not candidates:
        return raw.decode("utf-8", errors="ignore")
    return max(candidates, key=_text_quality_score)


def _normalize_pdf_extracted_text(text: str) -> str:
    if not text:
        return ""
    normalized = text.replace(chr(0), "")
    normalized = normalized.replace("\u3000", " ")
    normalized = re.sub(r"[ \t]+\n", "\n", normalized)
    normalized = re.sub(r"\n{3,}", "\n\n", normalized)
    return normalized.strip()


def _pdf_text_quality_issue(text: str, pages_total: int = 0, pages_with_text: int = 0) -> str:
    normalized = _normalize_pdf_extracted_text(text)
    char_count = len(normalized.strip())
    if char_count <= 0:
        return "PDF \u672a\u63d0\u53d6\u5230\u53ef\u7528\u6587\u672c"

    lines = [line.strip() for line in normalized.splitlines() if line.strip()]
    url_pattern = re.compile(r"(?:https?://\S+|www\.\S+)", flags=re.I)
    page_url_only_lines = sum(
        1
        for line in lines
        if re.fullmatch(r"\[page\s+\d+\]\s*(?:(?:https?://\S+|www\.\S+)\s*)+", line, flags=re.I)
    )
    url_hits = len(url_pattern.findall(normalized))

    semantic_probe = re.sub(r"\[page\s+\d+\]", " ", normalized, flags=re.I)
    semantic_probe = url_pattern.sub(" ", semantic_probe)
    semantic_probe = re.sub(r"[_=~\-]{3,}", " ", semantic_probe)
    semantic_probe = re.sub(r"\s+", " ", semantic_probe).strip()

    semantic_cjk = len(re.findall(r"[\u4e00-\u9fff]", semantic_probe))
    semantic_words = len(re.findall(r"[A-Za-z]{3,}", semantic_probe))
    semantic_score = semantic_cjk + semantic_words * 2

    if lines and page_url_only_lines >= max(6, int(len(lines) * 0.45)) and semantic_score < 60:
        return "\u63d0\u53d6\u7ed3\u679c\u4ee5\u91cd\u590d\u7f51\u5740/\u6c34\u5370\u4e3a\u4e3b"
    if url_hits >= max(10, pages_with_text or 10) and semantic_score < 80:
        return "\u63d0\u53d6\u7ed3\u679c\u4ee5\u91cd\u590d\u7f51\u5740/\u6c34\u5370\u4e3a\u4e3b"

    if pages_total <= 0:
        return "" if char_count >= 120 else "PDF \u6587\u672c\u8fc7\u77ed"

    coverage = pages_with_text / max(1, pages_total)
    if char_count >= 800:
        return ""
    if coverage >= 0.35 and char_count >= 180:
        return ""
    if pages_total <= 3 and pages_with_text >= 1 and char_count >= 120:
        return ""
    return "PDF \u672a\u63d0\u53d6\u5230\u8db3\u591f\u53ef\u7528\u6587\u672c"

def _image_to_data_url(image_bytes: bytes, suffix: str) -> str:
    mime_map = {
        "png": "image/png",
        "jpg": "image/jpeg",
        "jpeg": "image/jpeg",
        "webp": "image/webp",
    }
    mime = mime_map.get(suffix.lower(), "image/png")
    encoded = base64.b64encode(image_bytes).decode("utf-8")
    return f"data:{mime};base64,{encoded}"


def _extract_text_with_vision(image_bytes: bytes, suffix: str, file_name: str) -> str:
    """Use OpenAI-compatible multimodal endpoint (e.g. DashScope compatible mode) to parse images."""
    try:
        from openai import OpenAI
    except Exception:
        return f"【多模态未启用】缺少 openai 包，无法解析图片：{file_name}"

    api_key = os.getenv("VISION_API_KEY") or os.getenv("DASHSCOPE_API_KEY") or os.getenv("OPENAI_API_KEY")
    base_url = os.getenv("VISION_API_BASE") or os.getenv("OPENAI_API_BASE")
    model = os.getenv("VISION_MODEL") or "qwen-vl-plus"
    if not api_key or not base_url:
        return f"【多模态未配置】请设置 VISION_API_KEY 与 VISION_API_BASE（或兼容的 OPENAI_API_*）后再解析图片：{file_name}"

    data_url = _image_to_data_url(image_bytes, suffix)
    prompt = (
        "你是课堂资料解析助手。请提取图片中的教学文本、题目、知识点与结构，"
        "输出为清晰的中文 Markdown，尽量保留原始逻辑层级。"
    )
    try:
        client = OpenAI(api_key=api_key, base_url=base_url)
        resp = client.chat.completions.create(
            model=model,
            temperature=0.2,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": data_url}},
                    ],
                }
            ],
        )
        content = ""
        if resp and getattr(resp, "choices", None):
            content = (resp.choices[0].message.content or "").strip()
        return content if content else f"【多模态空响应】{file_name}"
    except Exception as e:
        return f"【多模态解析失败】{file_name}: {str(e)}"


def process_image_with_vision_prompt(image_bytes: bytes, prompt: str, suffix: str = "png") -> str:
    """Run a custom multimodal prompt for an image using OpenAI-compatible API."""
    try:
        from openai import OpenAI
    except Exception as exc:
        raise RuntimeError(f"缺少 openai 包，无法调用视觉模型：{exc}") from exc

    api_key = os.getenv("VISION_API_KEY") or os.getenv("DASHSCOPE_API_KEY") or os.getenv("OPENAI_API_KEY")
    base_url = os.getenv("VISION_API_BASE") or os.getenv("OPENAI_API_BASE")
    model = os.getenv("VISION_MODEL") or "qwen-vl-plus"
    if not api_key or not base_url:
        raise RuntimeError("未配置视觉模型环境变量，请设置 VISION_API_KEY 与 VISION_API_BASE。")

    data_url = _image_to_data_url(image_bytes, suffix)
    client = OpenAI(api_key=api_key, base_url=base_url)
    resp = client.chat.completions.create(
        model=model,
        temperature=0.3,
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": data_url}},
                ],
            }
        ],
    )
    if not resp or not getattr(resp, "choices", None):
        return "【多模态空响应】"
    return (resp.choices[0].message.content or "").strip()


def _guess_uploaded_image_suffix(uploaded_obj) -> str:
    name = str(getattr(uploaded_obj, "name", "") or "").lower()
    if "." in name:
        ext = name.rsplit(".", 1)[-1]
        if ext in {"png", "jpg", "jpeg", "webp"}:
            return ext
    mime = str(getattr(uploaded_obj, "type", "") or "").lower()
    if "jpeg" in mime or "jpg" in mime:
        return "jpg"
    if "webp" in mime:
        return "webp"
    return "png"


def _is_probable_extraction_error_text(text: str) -> bool:
    normalized = (text or "").strip()
    if not normalized:
        return True
    return bool(
        re.match(
            r"^【(多模态未启用|多模态未配置|多模态解析失败|多模态空响应|解析失败|提示|不支持的文件类型)",
            normalized,
        )
    )


def _pdf_text_is_usable(text: str, pages_total: int, pages_with_text: int) -> bool:
    return not _pdf_text_quality_issue(text, pages_total, pages_with_text)

def _extract_pdf_text_with_pypdf(raw: bytes) -> tuple[str, dict]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(raw))
    pages: list[str] = []
    pages_total = len(reader.pages)
    pages_with_text = 0
    for idx, page in enumerate(reader.pages):
        text = (page.extract_text() or "").strip()
        if text:
            pages_with_text += 1
            pages.append(f"[Page {idx + 1}] {text}")
    extracted = _normalize_pdf_extracted_text("\n".join(pages).strip())
    return extracted, {
        "pages_total": pages_total,
        "pages_with_text": pages_with_text,
        "chars": len(extracted),
    }


def _extract_pdf_text_with_pdfplumber(raw: bytes) -> tuple[str, dict]:
    import pdfplumber

    pages: list[str] = []
    with pdfplumber.open(io.BytesIO(raw)) as pdf:
        pages_total = len(pdf.pages)
        pages_with_text = 0
        for idx, page in enumerate(pdf.pages):
            text = (page.extract_text() or "").strip()
            if text:
                pages_with_text += 1
                pages.append(f"[Page {idx + 1}] {text}")
    extracted = _normalize_pdf_extracted_text("\n".join(pages).strip())
    return extracted, {
        "pages_total": pages_total,
        "pages_with_text": pages_with_text,
        "chars": len(extracted),
    }


def _extract_text_from_uploaded_file_with_status(uploaded_file) -> tuple[str, dict]:
    file_name = str(getattr(uploaded_file, "name", "") or "未命名文件")
    name_lower = file_name.lower()
    raw = uploaded_file.getvalue()
    suffix = name_lower.rsplit(".", 1)[-1] if "." in name_lower else ""
    status = {
        "file_name": file_name,
        "file_type": suffix or "unknown",
        "ok": False,
        "method": "",
        "chars": 0,
        "pages_total": None,
        "pages_with_text": None,
        "message": "",
        "detail": "",
    }

    try:
        if suffix in {"txt", "md", "csv", "json"}:
            text = (_decode_text_bytes(raw) or "").strip()
            status.update(
                {
                    "ok": bool(text),
                    "method": "plain_text",
                    "chars": len(text),
                    "message": "解析成功" if text else "文本为空",
                }
            )
            return text, status

        if suffix in {"html", "htm"}:
            html_text = _decode_text_bytes(raw)
            no_script = re.sub(r"(?is)<(script|style).*?>.*?</\\1>", " ", html_text)
            no_tag = re.sub(r"(?is)<[^>]+>", " ", no_script)
            text = html.unescape(re.sub(r"[ \t]+", " ", no_tag)).strip()
            status.update(
                {
                    "ok": bool(text),
                    "method": "html_strip",
                    "chars": len(text),
                    "message": "解析成功" if text else "HTML 去标签后为空",
                }
            )
            return text, status

        if suffix == "pdf":
            pypdf_text = ""
            pypdf_meta: dict = {}
            pypdf_error = ""
            try:
                pypdf_text, pypdf_meta = _extract_pdf_text_with_pypdf(raw)
            except Exception as exc:
                pypdf_error = str(exc)

            pypdf_pages_total = int(pypdf_meta.get("pages_total") or 0)
            pypdf_pages_with_text = int(pypdf_meta.get("pages_with_text") or 0)
            pypdf_usable = _pdf_text_is_usable(pypdf_text, pypdf_pages_total, pypdf_pages_with_text)

            chosen_text = (pypdf_text or "").strip()
            chosen_method = "pypdf"
            chosen_meta = pypdf_meta
            detail_parts: list[str] = []
            if pypdf_error:
                detail_parts.append(f"pypdf失败: {pypdf_error}")
            else:
                detail_parts.append(
                    f"pypdf页覆盖 {pypdf_pages_with_text}/{pypdf_pages_total}，文本 {len(chosen_text)} 字符"
                )

            if not pypdf_usable or pypdf_error:
                pdfplumber_text = ""
                pdfplumber_meta: dict = {}
                pdfplumber_error = ""
                try:
                    pdfplumber_text, pdfplumber_meta = _extract_pdf_text_with_pdfplumber(raw)
                except Exception as exc:
                    pdfplumber_error = str(exc)

                pdfplumber_pages_total = int(pdfplumber_meta.get("pages_total") or 0)
                pdfplumber_pages_with_text = int(pdfplumber_meta.get("pages_with_text") or 0)
                pdfplumber_usable = _pdf_text_is_usable(
                    pdfplumber_text, pdfplumber_pages_total, pdfplumber_pages_with_text
                )

                if pdfplumber_error:
                    detail_parts.append(f"pdfplumber失败: {pdfplumber_error}")
                else:
                    detail_parts.append(
                        f"pdfplumber页覆盖 {pdfplumber_pages_with_text}/{pdfplumber_pages_total}，文本 {len(pdfplumber_text.strip())} 字符"
                    )

                if pdfplumber_usable and (not pypdf_usable or len(pdfplumber_text.strip()) >= len(chosen_text)):
                    chosen_text = (pdfplumber_text or "").strip()
                    chosen_method = "pypdf -> pdfplumber"
                    chosen_meta = pdfplumber_meta
                elif not chosen_text and (pdfplumber_text or "").strip():
                    chosen_text = (pdfplumber_text or "").strip()
                    chosen_method = "pdfplumber(低置信)"
                    chosen_meta = pdfplumber_meta

            if not chosen_text:
                status.update(
                    {
                        "ok": False,
                        "method": chosen_method,
                        "message": "PDF 未提取到可用文本",
                        "detail": "；".join(detail_parts),
                    }
                )
                return "", status

            quality_issue = _pdf_text_quality_issue(
                chosen_text,
                int(chosen_meta.get("pages_total") or 0),
                int(chosen_meta.get("pages_with_text") or 0),
            )
            if quality_issue:
                status.update(
                    {
                        "ok": False,
                        "method": chosen_method,
                        "chars": len(chosen_text),
                        "pages_total": int(chosen_meta.get("pages_total") or 0) or None,
                        "pages_with_text": int(chosen_meta.get("pages_with_text") or 0) or None,
                        "message": quality_issue,
                        "detail": "；".join(detail_parts),
                    }
                )
                return "", status

            status.update(
                {
                    "ok": True,
                    "method": chosen_method,
                    "chars": len(chosen_text),
                    "pages_total": int(chosen_meta.get("pages_total") or 0) or None,
                    "pages_with_text": int(chosen_meta.get("pages_with_text") or 0) or None,
                    "message": "解析成功",
                    "detail": "；".join(detail_parts),
                }
            )
            return chosen_text, status


        if suffix == "docx":
            doc = Document(io.BytesIO(raw))
            text = "\n".join(p.text for p in doc.paragraphs if p.text and p.text.strip()).strip()
            status.update(
                {
                    "ok": bool(text),
                    "method": "docx_paragraphs",
                    "chars": len(text),
                    "message": "解析成功" if text else "DOCX 中未识别到段落文本",
                }
            )
            return text, status

        if suffix == "pptx":
            prs = Presentation(io.BytesIO(raw))
            lines: list[str] = []
            for slide_idx, slide in enumerate(prs.slides, start=1):
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text and text.strip():
                        lines.append(f"[Slide {slide_idx}] {text.strip()}")
            parsed = "\n".join(lines).strip()
            status.update(
                {
                    "ok": bool(parsed),
                    "method": "pptx_shapes",
                    "chars": len(parsed),
                    "message": "解析成功" if parsed else "PPT 中未识别到可提取文本",
                }
            )
            return parsed, status

        if suffix in {"png", "jpg", "jpeg", "webp"}:
            vision_text = (_extract_text_with_vision(raw, suffix, file_name) or "").strip()
            if _is_probable_extraction_error_text(vision_text):
                status.update(
                    {
                        "ok": False,
                        "method": "vision_api",
                        "message": vision_text or "视觉模型未返回文本",
                    }
                )
                return "", status

            status.update(
                {
                    "ok": True,
                    "method": "vision_api",
                    "chars": len(vision_text),
                    "message": "解析成功",
                }
            )
            return vision_text, status

        status.update({"ok": False, "method": "unknown", "message": f"不支持的文件类型: {file_name}"})
        return "", status
    except Exception as exc:
        status.update({"ok": False, "message": f"解析失败: {str(exc)}"})
        return "", status


def _extract_text_from_uploaded_file(uploaded_file) -> str:
    extracted, _ = _extract_text_from_uploaded_file_with_status(uploaded_file)
    return extracted


def _split_text_windows(text: str, chunk_size: int = 1800, overlap: int = 320) -> list[str]:
    cleaned = re.sub(r"\n{3,}", "\n\n", text or "").strip()
    if not cleaned:
        return []
    windows: list[str] = []
    step = max(1, chunk_size - overlap)
    for start in range(0, len(cleaned), step):
        window = cleaned[start:start + chunk_size].strip()
        if window:
            windows.append(window)
    return windows


def _select_relevant_snippets(text: str, query_text: str, max_chars: int = 6500, max_snippets: int = 4) -> tuple[str, str]:
    cleaned = (text or "").strip()
    query = (query_text or "").strip()
    if not cleaned:
        return "", "empty"
    if not query:
        return cleaned[:max_chars], "front"

    windows = _split_text_windows(cleaned)
    if not windows:
        return cleaned[:max_chars], "front"

    scored: list[tuple[int, int, str]] = []
    for idx, window in enumerate(windows):
        score = _score_text_match(query, window)
        if score > 0:
            scored.append((score, idx, window))

    if not scored:
        terms = sorted(_build_query_terms(query), key=len, reverse=True)
        lower = cleaned.lower()
        for term in terms:
            pos = lower.find(term.lower())
            if pos >= 0:
                start = max(0, pos - max_chars // 3)
                end = min(len(cleaned), start + max_chars)
                return cleaned[start:end].strip(), f"keyword_window:{term}"
        return cleaned[:max_chars], "front_no_match"

    scored.sort(key=lambda item: item[0], reverse=True)
    selected: list[tuple[int, int, str]] = []
    for score, idx, window in scored:
        if any(abs(idx - existing_idx) <= 1 for _, existing_idx, _ in selected):
            continue
        selected.append((score, idx, window))
        if len(selected) >= max_snippets:
            break

    if not selected:
        selected = scored[:max_snippets]

    selected.sort(key=lambda item: item[1])
    parts: list[str] = []
    total = 0
    for rank, (score, idx, window) in enumerate(selected, start=1):
        header = f"[相关片段 {rank} | 窗口 {idx + 1} | 匹配分 {score}]"
        block = f"{header}\n{window}"
        if total + len(block) > max_chars and parts:
            break
        parts.append(block)
        total += len(block)

    return "\n\n".join(parts).strip()[:max_chars], f"topic_ranked:{len(parts)}"


def build_external_knowledge_context(uploaded_files: list, query_text: str = "") -> tuple[str, list[str], list[dict]]:
    if not uploaded_files:
        return "", [], []

    source_names: list[str] = []
    chunks: list[str] = []
    parse_statuses: list[dict] = []
    max_chars_per_file = 6500
    max_total_chars = 22000

    for file in uploaded_files:
        source_names.append(file.name)
        extracted, status = _extract_text_from_uploaded_file_with_status(file)
        parse_statuses.append(status)
        if not extracted.strip():
            continue

        snippet, select_method = _select_relevant_snippets(
            extracted,
            query_text=query_text,
            max_chars=max_chars_per_file,
        )
        if query_text and status.get("ok"):
            detail = str(status.get("detail") or "").strip()
            selection_detail = f"主题片段选择: {select_method}"
            status["detail"] = f"{detail}；{selection_detail}" if detail else selection_detail
        chunks.append(f"[资料来源: {file.name}]\n{snippet}")

        if sum(len(c) for c in chunks) >= max_total_chars:
            break

    context = "\n\n".join(chunks)
    if len(context) > max_total_chars:
        context = context[:max_total_chars]
    return context, source_names, parse_statuses


def _kb_storage_dir() -> str:
    path = os.path.join(os.getcwd(), ".teacher_kb")
    os.makedirs(path, exist_ok=True)
    return path


def _fallback_kb_path() -> str:
    return os.path.join(_kb_storage_dir(), "knowledge_chunks.json")


def _load_fallback_kb() -> list[dict]:
    path = _fallback_kb_path()
    if not os.path.exists(path):
        return []
    try:
        with open(path, "r", encoding="utf-8") as f:
            data = json.load(f)
        if not isinstance(data, list):
            return []
        repaired: list[dict] = []
        for item in data:
            if not isinstance(item, dict):
                continue
            normalized = dict(item)
            normalized["source"] = _repair_mojibake_text(str(normalized.get("source", "")))
            normalized_text = _repair_mojibake_text(str(normalized.get("text", "")))
            if str(normalized["source"]).lower().endswith(".pdf"):
                normalized_text = _normalize_pdf_extracted_text(normalized_text)
                if _pdf_text_quality_issue(normalized_text):
                    continue
            normalized["text"] = normalized_text
            if not str(normalized.get("text", "") or "").strip():
                continue
            repaired.append(normalized)
        return repaired
    except Exception:
        return []

def _save_fallback_kb(records: list[dict]) -> None:
    path = _fallback_kb_path()
    with open(path, "w", encoding="utf-8") as f:
        json.dump(records, f, ensure_ascii=False)


def _dedup_kb_records(records: list[dict]) -> list[dict]:
    dedup_map: dict[str, dict] = {}
    for item in records:
        if not isinstance(item, dict):
            continue
        doc_id = str(item.get("id") or "").strip()
        if not doc_id:
            continue
        dedup_map[doc_id] = item
    return list(dedup_map.values())


def _split_for_kb(text: str, chunk_size: int = 1100, overlap: int = 180) -> list[str]:
    cleaned = re.sub(r"\n{3,}", "\n\n", text or "").strip()
    if not cleaned:
        return []
    chunks: list[str] = []
    step = max(1, chunk_size - overlap)
    for start in range(0, len(cleaned), step):
        chunk = cleaned[start:start + chunk_size].strip()
        if chunk:
            chunks.append(chunk)
    return chunks


def _get_chroma_collection():
    if not HAS_CHROMA:
        return None
    try:
        db_path = os.path.join(_kb_storage_dir(), "chroma_db")
        client = chromadb.PersistentClient(path=db_path)
        return client.get_or_create_collection(name="teacher_materials")
    except Exception:
        return None


def add_files_to_teacher_kb(uploaded_files: list) -> tuple[int, int, list[dict]]:
    """Return (files_added, chunks_added, parse_statuses)."""
    if not uploaded_files:
        return 0, 0, []

    files_added = 0
    chunks_added = 0
    parse_statuses: list[dict] = []
    collection = _get_chroma_collection()
    fallback_records = _load_fallback_kb()
    newly_built_records: list[dict] = []

    for file in uploaded_files:
        extracted, status = _extract_text_from_uploaded_file_with_status(file)
        parse_statuses.append(status)
        extracted = extracted.strip()
        if not extracted:
            continue

        chunks = _split_for_kb(extracted)
        if not chunks:
            continue

        files_added += 1
        for idx, chunk in enumerate(chunks):
            digest = hashlib.md5(f"{file.name}:{idx}:{chunk[:120]}".encode("utf-8", errors="ignore")).hexdigest()
            doc_id = f"kb_{digest}"
            metadata = {"source": file.name, "index": idx}
            record = {
                "id": doc_id,
                "source": file.name,
                "text": chunk,
            }
            newly_built_records.append(record)

            if collection is not None:
                try:
                    collection.add(
                        documents=[chunk],
                        ids=[doc_id],
                        metadatas=[metadata],
                    )
                except Exception:
                    # id 冲突时忽略，避免重复入库报错
                    pass
            chunks_added += 1

    if newly_built_records:
        merged = _dedup_kb_records([*fallback_records, *newly_built_records])
        _save_fallback_kb(merged)

    return files_added, chunks_added, parse_statuses


def add_text_to_teacher_kb(raw_text: str, source_name: str = "智能作答与批改引擎") -> int:
    text = (raw_text or "").strip()
    if not text:
        return 0

    chunks = _split_for_kb(text)
    if not chunks:
        return 0

    collection = _get_chroma_collection()
    fallback_records = _load_fallback_kb()
    new_records: list[dict] = []

    for idx, chunk in enumerate(chunks):
        digest = hashlib.md5(f"{source_name}:{idx}:{chunk[:120]}".encode("utf-8", errors="ignore")).hexdigest()
        doc_id = f"kb_{digest}"
        metadata = {"source": source_name, "index": idx}
        new_records.append({"id": doc_id, "source": source_name, "text": chunk})

        if collection is not None:
            try:
                collection.add(
                    documents=[chunk],
                    ids=[doc_id],
                    metadatas=[metadata],
                )
            except Exception:
                pass

    merged = _dedup_kb_records([*fallback_records, *new_records])
    _save_fallback_kb(merged)
    return len(chunks)


def _build_query_terms(query: str) -> list[str]:
    q = (query or "").lower().strip()
    if not q:
        return []

    stop_terms = {
        "解析",
        "分析",
        "相关",
        "章节",
        "讲义",
        "课程",
        "教学",
        "生成",
        "基于",
        "上传",
        "资料",
        "内容",
        "进行",
        "要求",
        "人物形象",
        "人物",
        "形象",
        "特点",
        "讲解",
        "介绍",
        "说明",
        "探究",
    }

    expanded: list[str] = []
    for term in re.split(r"[\s,，。;；、|/]+", q):
        term = term.strip()
        has_cjk = bool(re.search(r"[\u4e00-\u9fff]", term))
        if len(term) >= 2 and term not in stop_terms and (not has_cjk or len(term) <= 8):
            expanded.append(term)

    if re.search(r"[\u4e00-\u9fff]", q):
        cjk_text = re.sub(r"[^\u4e00-\u9fff]+", " ", q)
        stop_pattern = "|".join(re.escape(term) for term in sorted(stop_terms, key=len, reverse=True))
        cjk_parts = [part.strip() for part in re.split(stop_pattern, cjk_text) if part.strip()]
        for part in cjk_parts:
            if 2 <= len(part) <= 8:
                expanded.append(part)
            for size in range(min(6, len(part)), 1, -1):
                for idx in range(0, max(0, len(part) - size + 1)):
                    term = part[idx:idx + size]
                    if term in stop_terms:
                        continue
                    if any(stop in term and len(term) <= len(stop) + 1 for stop in stop_terms):
                        continue
                    expanded.append(term)

    if re.search(r"[A-Za-z0-9]", q):
        for term in re.findall(r"[A-Za-z][A-Za-z0-9_.+-]{1,}", q):
            expanded.append(term.lower())

    ordered: list[str] = []
    for term in expanded:
        if len(term) < 2:
            continue
        if term not in ordered:
            ordered.append(term)
        if len(ordered) >= 160:
            break
    return ordered


def _query_term_weights(query: str) -> dict[str, float]:
    q = (query or "").lower().strip()
    weights: dict[str, float] = {}
    for term in _build_query_terms(q):
        query_hits = max(1, q.count(term))
        length_weight = min(10.0, float(len(term))) ** 1.35
        weights[term] = max(weights.get(term, 0.0), query_hits * length_weight)
    return weights


def _important_query_terms(query: str, limit: int = 12) -> list[str]:
    weights = _query_term_weights(query)
    candidates = [
        (term, weight)
        for term, weight in weights.items()
        if 3 <= len(term) <= 5 and not re.fullmatch(r"[a-z]{1,2}", term, flags=re.I)
    ]
    query_low = (query or "").lower()
    def importance_key(item: tuple[str, float]) -> tuple[int, int, float]:
        term, weight = item
        length_preference = {4: 4, 3: 3, 5: 2}.get(len(term), 1)
        return query_low.count(term), length_preference, weight

    candidates.sort(key=importance_key, reverse=True)
    chosen: list[str] = []
    for term, _ in candidates:
        if any(term in existing or existing in term for existing in chosen):
            if not any(term == existing for existing in chosen):
                continue
        chosen.append(term)
        if len(chosen) >= limit:
            break
    return chosen


def _score_text_match(query: str, text: str) -> int:
    query_low = (query or "").lower().strip()
    text_low = (text or "").lower()
    if not query_low or not text_low:
        return 0

    weights = _query_term_weights(query)
    if not weights:
        return 0

    matched_terms = 0
    raw_hits = 0
    weighted_hits = 0.0
    for term, weight in weights.items():
        cnt = text_low.count(term)
        if cnt > 0:
            matched_terms += 1
            raw_hits += cnt
            weighted_hits += weight * min(cnt, 4)

    if matched_terms == 0:
        return 0

    important_terms = _important_query_terms(query)
    important_total = sum(weights.get(term, 0.0) for term in important_terms) or 1.0
    important_hit = sum(weights.get(term, 0.0) for term in important_terms if term in text_low)
    important_coverage = important_hit / important_total
    coverage = matched_terms / max(1, min(len(weights), 32))
    score = int(
        weighted_hits * 24
        + raw_hits * 4
        + matched_terms * 10
        + min(1.0, coverage) * 420
        + min(1.0, important_coverage) * 900
    )

    if len(weights) >= 6 and matched_terms <= 1:
        score -= 180
    return max(0, score)


def _extract_preview_around_query(text: str, query: str, width: int = 220) -> str:
    raw = (text or "").strip()
    if not raw:
        return ""
    compact = re.sub(r"\s+", " ", raw)
    terms = sorted(_build_query_terms(query), key=len, reverse=True)
    center_idx = -1
    lower = compact.lower()
    for term in terms:
        idx = lower.find(term)
        if idx >= 0:
            center_idx = idx
            break
    if center_idx < 0:
        return compact[:width]
    start = max(0, center_idx - width // 3)
    end = min(len(compact), start + width)
    return compact[start:end]


def query_teacher_kb_hits(query_text: str, n_results: int = 4) -> list[dict]:
    chroma_hits: list[dict] = []
    collection = _get_chroma_collection()
    if collection is not None:
        try:
            if collection.count() > 0:
                results = collection.query(query_texts=[query_text], n_results=n_results)
                docs = results.get("documents", [[]])[0] if results else []
                metas = results.get("metadatas", [[]])[0] if results else []
                for i, doc in enumerate(docs):
                    txt = (doc or "").strip()
                    if not txt:
                        continue
                    source = "未知来源"
                    if i < len(metas) and isinstance(metas[i], dict):
                        source = str(metas[i].get("source", "未知来源"))
                    keyword_score = _score_text_match(query_text, txt)
                    chroma_hits.append(
                        {
                            "source": source,
                            "text": txt,
                            "score": int(keyword_score) if keyword_score > 0 else None,
                        }
                    )
        except Exception:
            pass

    hits: list[dict] = []
    records = _load_fallback_kb()
    if records:
        term_weights = _query_term_weights(query_text)
        doc_freq: dict[str, int] = {}
        for term in term_weights:
            doc_freq[term] = sum(1 for item in records if term in (item.get("text", "") or "").lower())

        def local_kb_score(item: dict) -> int:
            text_value = item.get("text", "") or ""
            text_low = text_value.lower()
            base_score = _score_text_match(query_text, text_value)
            idf_score = 0.0
            for term, weight in term_weights.items():
                cnt = text_low.count(term)
                if cnt <= 0:
                    continue
                idf = math.log((len(records) + 1) / (doc_freq.get(term, 0) + 1)) + 1.0
                idf_score += weight * idf * min(cnt, 4)
            return int(base_score + idf_score * 18)

        ranked = sorted(records, key=local_kb_score, reverse=True)
        scored_top = []
        for rec in ranked[: max(n_results * 8, 30)]:
            score = local_kb_score(rec)
            if score > 0:
                scored_top.append((rec, score))

        if scored_top:
            scored_top.sort(key=lambda x: x[1], reverse=True)
            max_score = scored_top[0][1]
            keep_floor = max(140, int(max_score * 0.30))
            scored_top = [pair for pair in scored_top if pair[1] >= keep_floor][:n_results]

            for rec, score in scored_top:
                src = rec.get("source", "未知来源")
                txt = (rec.get("text", "") or "").strip()
                if txt:
                    hits.append(
                        {
                            "source": src,
                            "text": txt,
                            "score": int(score),
                        }
                    )

    seen: set[str] = set()
    merged: list[dict] = []
    for item in [*hits, *chroma_hits]:
        txt = (item.get("text", "") or "").strip()
        if not txt:
            continue
        fingerprint = hashlib.md5(txt[:300].encode("utf-8", errors="ignore")).hexdigest()
        if fingerprint in seen:
            continue
        seen.add(fingerprint)
        merged.append(item)

    if not merged:
        return []

    merged.sort(
        key=lambda item: item.get("score") if isinstance(item.get("score"), int) else -1,
        reverse=True,
    )
    return merged[:n_results]


def query_teacher_kb(query_text: str, n_results: int = 4) -> str:
    hits = query_teacher_kb_hits(query_text, n_results=n_results)
    if not hits:
        return ""
    lines = [f"[知识库:{hit.get('source', '未知来源')}] {(hit.get('text', '') or '').strip()}" for hit in hits if (hit.get("text") or "").strip()]
    return "\n\n".join(lines).strip()


def render_parse_status_panel(statuses: list[dict], title: str) -> None:
    rows = [s for s in (statuses or []) if isinstance(s, dict)]
    if not rows:
        return

    ok_count = sum(1 for s in rows if s.get("ok"))
    fail_count = len(rows) - ok_count
    st.caption(f"{title}：成功 {ok_count}，失败 {fail_count}")

    with st.expander(f"查看{title}详情", expanded=False):
        for idx, item in enumerate(rows, start=1):
            file_name = str(item.get("file_name", "未知文件"))
            result = "成功" if item.get("ok") else "失败"
            method = str(item.get("method", "-"))
            chars = int(item.get("chars") or 0)

            pages_total = item.get("pages_total")
            pages_with_text = item.get("pages_with_text")
            page_info = ""
            if isinstance(pages_total, int) and pages_total > 0:
                covered = pages_with_text if isinstance(pages_with_text, int) else 0
                page_info = f"，页文本覆盖 {covered}/{pages_total}"

            st.markdown(f"{idx}. {file_name} | {result} | 方法 {method} | 文本 {chars} 字符{page_info}")

            message = str(item.get("message", "") or "").strip()
            detail = str(item.get("detail", "") or "").strip()
            if message:
                st.caption(f"原因: {message}")
            if detail:
                st.caption(f"细节: {detail}")


def teacher_kb_chunk_count() -> int:
    fallback_count = len(_load_fallback_kb())
    collection = _get_chroma_collection()
    if collection is not None:
        try:
            chroma_count = int(collection.count())
            return max(chroma_count, fallback_count)
        except Exception:
            return fallback_count
    return fallback_count


def clear_teacher_kb() -> tuple[int, str]:
    previous_count = teacher_kb_chunk_count()
    storage_dir = _kb_storage_dir()
    fallback_path = _fallback_kb_path()
    chroma_dir = os.path.join(storage_dir, "chroma_db")

    errors: list[str] = []
    try:
        if os.path.exists(fallback_path):
            os.remove(fallback_path)
    except Exception as exc:
        errors.append(f"删除 knowledge_chunks.json 失败: {str(exc)}")

    try:
        if os.path.isdir(chroma_dir):
            shutil.rmtree(chroma_dir, ignore_errors=False)
    except Exception as exc:
        errors.append(f"删除 chroma_db 失败: {str(exc)}")

    if errors:
        return previous_count, "；".join(errors)
    return previous_count, "知识库已清空。"


def extract_python_code_blocks(md_text: str) -> list[str]:
    pattern = re.compile(r"```(?:python|py)\s*(.*?)```", flags=re.IGNORECASE | re.DOTALL)
    blocks = [block.strip() for block in pattern.findall(md_text or "")]
    return [block for block in blocks if block]


def _extract_missing_module(error_text: str) -> str | None:
    if not error_text:
        return None
    match = re.search(r"No module named ['\"]([^'\"]+)['\"]", error_text)
    if not match:
        return None
    module_name = match.group(1).strip()
    if re.fullmatch(r"[a-zA-Z0-9_.-]+", module_name):
        return module_name
    return None


def _install_missing_module(module_name: str, timeout_sec: int = 120) -> tuple[bool, str]:
    try:
        process = subprocess.run(
            [sys.executable, "-m", "pip", "install", module_name],
            capture_output=True,
            text=True,
            timeout=timeout_sec,
            check=False,
        )
        output = (process.stdout or "") + "\n" + (process.stderr or "")
        output = output.strip()
        return process.returncode == 0, output[-1800:] if output else ""
    except subprocess.TimeoutExpired:
        return False, f"安装模块 {module_name} 超时（>{timeout_sec}秒）。"
    except Exception as exc:
        return False, f"安装模块 {module_name} 失败: {exc}"


def _run_code_block(code: str, timeout_sec: int = 8) -> dict:
    start = time.perf_counter()
    try:
        process = subprocess.run(
            [sys.executable, "-c", code],
            capture_output=True,
            text=True,
            timeout=timeout_sec,
            check=False,
        )
        duration_ms = int((time.perf_counter() - start) * 1000)
        stdout = (process.stdout or "").strip()
        stderr = (process.stderr or "").strip()
        success = process.returncode == 0

        if success:
            output = stdout if stdout else "代码执行成功，无标准输出。"
        else:
            output = stderr or stdout or "代码执行失败，但未返回详细错误信息。"
            missing_module = _extract_missing_module(output)
            if missing_module:
                installed, install_log = _install_missing_module(missing_module)
                if installed:
                    retry = subprocess.run(
                        [sys.executable, "-c", code],
                        capture_output=True,
                        text=True,
                        timeout=timeout_sec,
                        check=False,
                    )
                    retry_success = retry.returncode == 0
                    retry_stdout = (retry.stdout or "").strip()
                    retry_stderr = (retry.stderr or "").strip()
                    if retry_success:
                        output = (
                            f"检测到缺失模块 {missing_module}，已自动安装并重试成功。\n"
                            f"重试输出:\n{retry_stdout if retry_stdout else '代码执行成功，无标准输出。'}"
                        )
                        success = True
                    else:
                        output = (
                            f"检测到缺失模块 {missing_module}，已自动安装但重试仍失败。\n"
                            f"重试错误:\n{retry_stderr or retry_stdout or '无详细错误'}"
                        )
                else:
                    output = (
                        f"检测到缺失模块 {missing_module}，已尝试自动安装但失败。\n"
                        f"安装日志（截断）:\n{install_log if install_log else '无日志'}\n"
                        f"原始错误:\n{stderr or stdout or '无详细错误'}"
                    )

        return {
            "success": success,
            "output": output,
            "duration_ms": duration_ms,
        }
    except subprocess.TimeoutExpired:
        duration_ms = int((time.perf_counter() - start) * 1000)
        return {
            "success": False,
            "output": f"执行超时（>{timeout_sec}秒），请检查是否存在死循环或阻塞操作。",
            "duration_ms": duration_ms,
        }
    except Exception as exc:
        duration_ms = int((time.perf_counter() - start) * 1000)
        return {
            "success": False,
            "output": f"执行器异常: {exc}",
            "duration_ms": duration_ms,
        }

def _extract_imported_modules(code: str) -> list[str]:
    modules = set()
    try:
        tree = ast.parse(code)
        for node in ast.walk(tree):
            if isinstance(node, ast.Import):
                for alias in node.names:
                    modules.add(alias.name.split('.')[0])
            elif isinstance(node, ast.ImportFrom):
                if node.module:
                    modules.add(node.module.split('.')[0])
    except Exception:
        pass
    
    if hasattr(sys, "stdlib_module_names"):
        # filter out stdlib and pseudo-modules
        modules = {m for m in modules if m not in sys.stdlib_module_names and m not in ("__future__",)}
    
    return sorted(list(modules))

def prepare_code_runner_agent(md_text: str) -> dict:
    blocks = extract_python_code_blocks(md_text)
    reports = []
    for idx, code in enumerate(blocks, start=1):
        imported_modules = _extract_imported_modules(code)
        reports.append(
            {
                "index": idx,
                "code": code,
                "success": False,
                "is_pending": True,
                "output": "等待用户确认执行...",
                "duration_ms": 0,
                "imported_modules": imported_modules,
            }
        )
    return {
        "status": "pending",
        "reports": reports,
        "total": len(reports),
        "passed": 0,
        "failed": 0,
        "total_duration_ms": 0,
    }

def run_code_runner_agent(md_text: str) -> dict:
    blocks = extract_python_code_blocks(md_text)
    reports = []

    for idx, code in enumerate(blocks, start=1):
        result = _run_code_block(code)
        reports.append(
            {
                "index": idx,
                "code": code,
                "success": result["success"],
                "is_pending": False,
                "output": result["output"],
                "duration_ms": result["duration_ms"],
            }
        )

    passed = sum(1 for item in reports if item["success"])
    total_duration = sum(item["duration_ms"] for item in reports)
    return {
        "status": "completed",
        "total": len(reports),
        "passed": passed,
        "failed": len(reports) - passed,
        "total_duration_ms": total_duration,
        "reports": reports,
        "generated_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
    }


def build_code_audit_markdown(topic: str, audit: dict) -> str:
    header = [
        f"# {topic} 代码运行报告",
        "",
        f"- 生成时间: {audit.get('generated_at', '')}",
        f"- 代码块总数: {audit.get('total', 0)}",
        f"- 通过: {audit.get('passed', 0)}",
        f"- 失败: {audit.get('failed', 0)}",
        f"- 总耗时: {audit.get('total_duration_ms', 0)} ms",
        "",
    ]

    body: list[str] = []
    for item in audit.get("reports", []):
        status = "PASS" if item.get("success") else "FAIL"
        body.extend(
            [
                f"## 代码块 {item.get('index')} [{status}]",
                "```python",
                item.get("code", ""),
                "```",
                "输出：",
                "```text",
                item.get("output", ""),
                "```",
                "",
            ]
        )

    return "\n".join(header + body)

# ==================== 文件生成逻辑 ====================

def generate_docx(md_text: str) -> bytes:
    doc = Document()
    normal_style = doc.styles["Normal"]
    normal_style.font.name = "Microsoft YaHei"
    normal_style.font.size = DocxPt(11)
    doc.add_heading("教学讲义", 0)
    lines = md_text.split("\n")
    in_code_block = False
    for raw_line in lines:
        line = raw_line.rstrip()
        if line.startswith("```"):
            in_code_block = not in_code_block
            continue
        if in_code_block:
            paragraph = doc.add_paragraph()
            run = paragraph.add_run(line)
            run.font.name = "Consolas"
            run.font.size = DocxPt(10)
            continue
        stripped = line.strip()
        if not stripped: continue
        heading_match = re.match(r"^(#{1,6})\s+(.*)", stripped)
        if heading_match:
            level = min(len(heading_match.group(1)), 4)
            doc.add_heading(heading_match.group(2).replace("**", "").replace("__", ""), level=level)
            continue
            
        is_bullet = stripped.startswith("- ")
        is_numbered = bool(re.match(r"^\d+\.\s+", stripped))
        
        if is_bullet:
            paragraph = doc.add_paragraph(style="List Bullet")
            content = stripped[2:]
        elif is_numbered:
            paragraph = doc.add_paragraph(style="List Number")
            content = re.sub(r"^\d+\.\s+", "", stripped)
        else:
            paragraph = doc.add_paragraph()
            content = stripped
            
        if "**" in content:
            parts = content.split("**")
            for index, part in enumerate(parts):
                run = paragraph.add_run(part)
                if index % 2 == 1: run.bold = True
        else:
            paragraph.add_run(content)
    bio = io.BytesIO()
    doc.save(bio)
    return bio.getvalue()

def generate_pdf(md_text: str) -> bytes | None:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
        from reportlab.platypus import Paragraph, Preformatted, SimpleDocTemplate, Spacer
    except ImportError:
        return None

    try:
        pdfmetrics.getFont("STSong-Light")
    except KeyError:
        pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))

    buffer = io.BytesIO()
    document = SimpleDocTemplate(buffer, pagesize=A4, leftMargin=40, rightMargin=40, topMargin=36, bottomMargin=36)
    styles = getSampleStyleSheet()

    normal = ParagraphStyle("NormalCN", parent=styles["Normal"], fontName="STSong-Light", fontSize=11, leading=16)
    heading1 = ParagraphStyle("Heading1CN", parent=styles["Heading1"], fontName="STSong-Light", fontSize=18, leading=24)
    heading2 = ParagraphStyle("Heading2CN", parent=styles["Heading2"], fontName="STSong-Light", fontSize=14, leading=20)
    heading3 = ParagraphStyle("Heading3CN", parent=styles["Heading3"], fontName="STSong-Light", fontSize=12, leading=18)
    code = ParagraphStyle("CodeCN", parent=styles["Code"], fontName="STSong-Light", fontSize=9, leading=13)

    story = []
    in_code_block = False
    code_lines = []

    for raw_line in md_text.splitlines():
        line = raw_line.rstrip()
        if line.startswith("```"):
            if in_code_block:
                story.append(Preformatted("\n".join(code_lines), code))
                story.append(Spacer(1, 8))
                code_lines = []
            in_code_block = not in_code_block
            continue
        if in_code_block:
            code_lines.append(line)
            continue
        stripped = line.strip()
        if not stripped:
            story.append(Spacer(1, 6))
            continue
            
        styled_text = html.escape(stripped)
        styled_text = re.sub(r'\*\*(.*?)\*\*', r'<b>\1</b>', styled_text)
        styled_text = re.sub(r'__(.*?)__', r'<b>\1</b>', styled_text)
        
        if styled_text.startswith("# "):
            story.append(Paragraph(styled_text[2:], heading1))
            continue
        if styled_text.startswith("## "):
            story.append(Paragraph(styled_text[3:], heading2))
            continue
        if styled_text.startswith("### "):
            story.append(Paragraph(styled_text[4:], heading3))
            continue
        story.append(Paragraph(styled_text, normal))

    if code_lines:
        story.append(Preformatted("\n".join(code_lines), code))

    document.build(story)
    return buffer.getvalue()

def _rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def _remove_margins(text_frame):
    text_frame.margin_left = Inches(0)
    text_frame.margin_right = Inches(0)
    text_frame.margin_top = Inches(0)
    text_frame.margin_bottom = Inches(0)


def _normalize_for_compare(text: str) -> str:
    text = text.strip()
    text = re.sub(r"^第[一二三四五六七八九十\d]+[、.．\s]*", "", text)
    return re.sub(r"[\s\-\*#：:，,。.!?！？；;（）()【】\[\]《》<>\d]", "", text).lower()


def _is_option_line(line: str) -> bool:
    stripped = line.strip()
    return bool(re.match(r"^(?:[-*]\s*)?[A-DＡ-Ｄ][\.、\)）]\s+", stripped))


def _is_question_stem(line: str) -> bool:
    return bool(re.match(r"^\d+\.\s+", line.strip()))


def _split_markdown_sections(md_text: str, topic: str) -> tuple[str, list[dict]]:
    title = (topic or "教学讲义").strip() or "教学讲义"
    sections: list[dict] = []
    current_title = "讲义内容"
    current_lines: list[str] = []
    in_code_block = False

    for raw_line in (md_text or "").splitlines():
        line = raw_line.rstrip()
        stripped = line.strip()

        if stripped.startswith("```"):
            current_lines.append("```")
            in_code_block = not in_code_block
            continue

        if not in_code_block and stripped.startswith("# "):
            maybe_title = stripped[2:].strip()
            if maybe_title:
                title = maybe_title
            continue

        if not in_code_block and stripped.startswith("## "):
            if current_lines:
                sections.append({"title": current_title, "lines": current_lines})
            current_title = stripped[3:].strip() or "讲义内容"
            current_lines = []
            continue

        if not stripped and not in_code_block:
            continue

        if in_code_block:
            current_lines.append(line)
        else:
            normalized = re.sub(r"^###\s+", "", stripped)
            if normalized:
                current_lines.append(normalized)

    if in_code_block:
        current_lines.append("```")

    if current_lines:
        sections.append({"title": current_title, "lines": current_lines})

    if not sections:
        sections = [{"title": "讲义内容", "lines": ["暂无可解析内容"]}]

    return title, sections


def _drop_repeated_leading_line(lines: list[str], section_title: str) -> list[str]:
    compact = [line for line in lines if line.strip()]
    if not compact:
        return []

    first = compact[0].strip()
    if first.startswith("```"):
        return compact

    first_norm = _normalize_for_compare(first)
    title_norm = _normalize_for_compare(section_title)
    if first_norm and title_norm and (first_norm == title_norm or first_norm in title_norm or title_norm in first_norm):
        return compact[1:]

    return compact


def _build_pagination_blocks(lines: list[str], max_code_lines: int = 8) -> list[list[str]]:
    compact = [line.rstrip() for line in lines if line.strip()]
    if not compact:
        return [["暂无可解析内容"]]

    blocks: list[list[str]] = []
    i = 0
    while i < len(compact):
        line = compact[i]

        if line.strip().startswith("```"):
            code_lines: list[str] = []
            i += 1
            while i < len(compact) and not compact[i].strip().startswith("```"):
                code_lines.append(compact[i])
                i += 1
            if i < len(compact) and compact[i].strip().startswith("```"):
                i += 1

            if not code_lines:
                blocks.append(["```", "# TODO: 补充示例代码", "pass", "```"])
                continue

            for start in range(0, len(code_lines), max_code_lines):
                segment = code_lines[start:start + max_code_lines]
                blocks.append(["```", *segment, "```"])
            continue

        if _is_question_stem(line):
            block = [line]
            i += 1
            while i < len(compact) and _is_option_line(compact[i]):
                block.append(compact[i])
                i += 1
            blocks.append(block)
            continue

        long_line_segments = _split_long_line_for_slide(line)
        for seg in long_line_segments:
            blocks.append([seg])
        i += 1

    return blocks


def _split_long_line_for_slide(line: str, max_len: int = 46) -> list[str]:
    stripped = line.strip()
    if not stripped or len(stripped) <= max_len:
        return [stripped] if stripped else []

    pieces: list[str] = []
    start = 0
    while start < len(stripped):
        end = min(len(stripped), start + max_len)
        window = stripped[start:end]
        punctuation = list(re.finditer(r"[，。；：,.!?！？]", window))
        if punctuation and len(window) > max_len // 2:
            cut = start + punctuation[-1].end()
        else:
            cut = end
        if cut <= start:
            cut = end
        segment = stripped[start:cut].strip()
        if segment:
            pieces.append(segment)
        start = cut
    return pieces or [stripped]


def _estimate_block_weight(block: list[str]) -> float:
    if block and block[0].strip().startswith("```"):
        inner_count = max(1, len(block) - 2)
        return 2.0 + inner_count * 0.9

    weight = 0.0
    for line in block:
        stripped = line.strip()
        base = 0.85 if _is_option_line(stripped) else 1.0
        length_weight = min(1.4, max(0.0, len(stripped) - 24) / 36.0)
        weight += base + length_weight
    return max(1.0, weight)


def _chunk_lines_for_slides(lines: list[str], max_weight: float = 14.0, max_code_lines: int = 14) -> list[list[str]]:
    # 代码块分页更宽松，避免被切得过碎
    blocks = _build_pagination_blocks(lines, max_code_lines=max_code_lines)
    if not blocks:
        return [["暂无可解析内容"]]

    chunks: list[list[str]] = []
    current_lines: list[str] = []
    current_weight = 0.0

    for block in blocks:
        block_weight = _estimate_block_weight(block)
        if current_lines and current_weight + block_weight > max_weight:
            chunks.append(current_lines)
            current_lines = block[:]
            current_weight = block_weight
            continue

        current_lines.extend(block)
        current_weight += block_weight

    if current_lines:
        chunks.append(current_lines)

    return chunks


def _clean_markdown_inline_for_slide(text: str) -> str:
    """Strip inline markdown markers that look noisy in plain-text PPT rendering."""
    if not text:
        return ""
    cleaned = text.replace("**", "").replace("__", "").replace("`", "")
    cleaned = cleaned.replace("*", "")
    return cleaned.strip()


def _ensure_run(paragraph):
    if paragraph.runs:
        return paragraph.runs[0]
    return paragraph.add_run()


def _add_cover_slide(prs: Presentation, topic: str, theme: str = "teaching_blue"):
    style = PPT_THEMES.get(theme, PPT_THEMES["teaching_blue"])
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(style["cover_bg"])

    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(style["cover_bg"])
    bg.line.fill.background()

    deco_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.84), Inches(1.12), Inches(0.24), Inches(5.2))
    deco_bar.fill.solid()
    deco_bar.fill.fore_color.rgb = _rgb(style["cover_deco"])
    deco_bar.line.fill.background()

    glow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.4), Inches(0.8), Inches(4.0), Inches(4.0))
    glow.fill.solid()
    glow.fill.fore_color.rgb = _rgb(style["cover_glow"])
    glow.fill.transparency = 0.45
    glow.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.34), Inches(1.52), Inches(10.9), Inches(2.3))
    _remove_margins(title_box.text_frame)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.text = topic
    title_para.alignment = PP_ALIGN.LEFT
    title_run = _ensure_run(title_para)
    title_run.font.name = "Microsoft YaHei"
    title_run.font.size = Pt(52)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(255, 255, 255)

    subtitle_box = slide.shapes.add_textbox(Inches(1.34), Inches(4.45), Inches(10.2), Inches(0.9))
    _remove_margins(subtitle_box.text_frame)
    subtitle_para = subtitle_box.text_frame.paragraphs[0]
    subtitle_para.text = "课堂讲义自动生成引擎 · Teaching Deck"
    subtitle_para.alignment = PP_ALIGN.LEFT
    subtitle_run = _ensure_run(subtitle_para)
    subtitle_run.font.name = "Microsoft YaHei"
    subtitle_run.font.size = Pt(19)
    subtitle_run.font.color.rgb = _rgb(style["cover_subtitle"])


def _add_content_slide(prs: Presentation, section_title: str, lines: list[str], continued: bool = False, theme: str = "teaching_blue"):
    style = PPT_THEMES.get(theme, PPT_THEMES["teaching_blue"])
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(style["page_bg"])

    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(style["page_bg"])
    bg.line.fill.background()

    header_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.08))
    header_band.fill.solid()
    header_band.fill.fore_color.rgb = _rgb(style["header_bg"])
    header_band.line.fill.background()

    content_card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(1.28), Inches(11.92), Inches(5.92))
    content_card.fill.solid()
    content_card.fill.fore_color.rgb = _rgb(style["card_bg"])
    content_card.line.color.rgb = _rgb(style["card_border"])
    content_card.line.width = Pt(1.4)

    full_title = section_title + ("（续）" if continued else "")
    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(0.24), Inches(11.6), Inches(0.75))
    _remove_margins(title_box.text_frame)
    title_para = title_box.text_frame.paragraphs[0]
    title_para.text = full_title
    title_para.alignment = PP_ALIGN.LEFT
    title_run = _ensure_run(title_para)
    title_run.font.name = "Microsoft YaHei"
    title_run.font.size = Pt(27)
    title_run.font.bold = True
    title_run.font.color.rgb = _rgb(style["title_text"])

    content_lines = _drop_repeated_leading_line(lines, section_title)
    if not content_lines:
        content_lines = ["暂无可解析内容"]

    visible_text_lines = [l for l in content_lines if l.strip() and not l.strip().startswith("```")]
    body_font_size = 19
    if len(visible_text_lines) >= 10:
        body_font_size = 17
    if len(visible_text_lines) >= 14:
        body_font_size = 15

    content_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.72), Inches(11.0), Inches(5.24))
    content_frame = content_box.text_frame
    _remove_margins(content_frame)
    content_frame.word_wrap = True
    # 左对齐不变，整体在内容卡片内垂直居中
    content_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    in_code_block = False
    has_written = False

    for raw_line in content_lines:
        stripped = raw_line.strip()
        if not stripped:
            continue

        if stripped.startswith("```"):
            in_code_block = not in_code_block
            continue

        paragraph = content_frame.paragraphs[0] if not has_written else content_frame.add_paragraph()
        has_written = True
        paragraph.alignment = PP_ALIGN.LEFT

        if in_code_block:
            paragraph.text = raw_line
            paragraph.space_after = Pt(2.5)
            code_run = _ensure_run(paragraph)
            code_run.font.name = "Consolas"
            code_run.font.size = Pt(13.2)
            code_run.font.color.rgb = _rgb(style["code_text"])
            continue

        text = re.sub(r"^###\s+", "", stripped)
        text = _clean_markdown_inline_for_slide(text)
        if text.startswith("- "):
            text = f"• {text[2:].strip()}"
        if not text:
            continue

        paragraph.text = text
        paragraph.space_after = Pt(5.6)
        if text.startswith("• "):
            paragraph.level = 0
        text_run = _ensure_run(paragraph)
        text_run.font.name = "Microsoft YaHei"
        text_run.font.size = Pt(16.2 if _is_option_line(text) else body_font_size)
        text_run.font.bold = _is_question_stem(text)
        text_run.font.color.rgb = _rgb(style["body_text"])


def generate_html_presentation(md_text: str, topic: str) -> bytes:
    cleaned_md = sanitize_markdown_for_ppt(md_text)
    _, sections = _split_markdown_sections(cleaned_md, topic)

    slides_md = [f"## {topic}\n\n自动生成教学讲义 · 浏览器演示版"]
    for section in sections:
        section_title = str(section.get("title") or "讲义内容")
        lines = section.get("lines") or ["暂无可解析内容"]
        chunks = _chunk_lines_for_slides(lines, max_weight=10.8, max_code_lines=9)

        for index, chunk in enumerate(chunks):
            title_line = f"### {section_title}{'（续）' if index > 0 else ''}"
            body_lines = _drop_repeated_leading_line(chunk, section_title)
            if not body_lines:
                body_lines = ["暂无可解析内容"]
            slides_md.append(f"{title_line}\n\n" + "\n".join(body_lines))

    markdown_content = "\n\n---\n\n".join(slides_md)
    markdown_content = markdown_content.replace("</textarea>", "&lt;/textarea&gt;")

    html_template = f"""<!DOCTYPE html>
<html lang=\"zh-CN\">
<head>
    <meta charset=\"utf-8\">
    <meta name=\"viewport\" content=\"width=device-width, initial-scale=1.0\">
    <title>{topic} - HTML 演示</title>
    <link rel=\"stylesheet\" href=\"https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reset.css\">
    <link rel=\"stylesheet\" href=\"https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reveal.css\">
    <link rel=\"stylesheet\" href=\"https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/theme/white.css\" id=\"theme\">
    <link rel=\"stylesheet\" href=\"https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/plugin/highlight/monokai.css\">
    <style>
        :root {{
            --bg-a: #edf2fb;
            --bg-b: #dbe7f6;
            --card: #ffffff;
            --ink: #1f2a44;
            --ink-soft: #43526a;
            --accent: #0f4c81;
            --accent-soft: #1f7a8c;
        }}
        .reveal {{
            font-family: 'Microsoft YaHei', 'PingFang SC', sans-serif;
            color: var(--ink);
            background:
                radial-gradient(45rem 28rem at 12% -8%, rgba(31, 122, 140, 0.17), transparent 72%),
                radial-gradient(40rem 26rem at 94% 6%, rgba(15, 76, 129, 0.18), transparent 74%),
                linear-gradient(145deg, var(--bg-a), var(--bg-b));
        }}
        .reveal .slides {{
            text-align: left;
        }}
        .reveal .slides section {{
            box-sizing: border-box;
            width: 100%;
            height: 100%;
            padding: 0.72rem 0.95rem !important;
            overflow-y: auto;
            overflow-x: hidden;
        }}
        .reveal .slides > section.present {{
            display: flex !important;
            align-items: stretch;
            justify-content: stretch;
        }}
        .reveal .slides > section > * {{
            width: 100%;
        }}
        .reveal h2 {{
            margin: 0 0 .5rem 0 !important;
            padding: .36rem .7rem !important;
            border-radius: 12px;
            background: linear-gradient(135deg, #123664, #1d4f85);
            color: #f6faff;
            font-size: 1.22em !important;
            letter-spacing: .01em;
            text-transform: none;
        }}
        .reveal h3 {{
            margin: 0.2rem 0 .45rem !important;
            color: var(--accent);
            font-size: 1.02em !important;
            text-transform: none;
            border-left: 4px solid var(--accent-soft);
            padding-left: .45rem;
        }}
        .reveal p, .reveal li {{
            color: var(--ink-soft);
            line-height: 1.56;
            font-size: .8em;
            word-break: break-word;
            overflow-wrap: anywhere;
        }}
        .reveal ul, .reveal ol {{
            margin: .25rem 0 .3rem 1rem;
        }}
        .reveal pre {{
            margin-top: .35rem;
            border-radius: 12px;
            box-shadow: 0 10px 24px rgba(27, 42, 68, 0.15);
            border: 1px solid rgba(31, 58, 95, 0.2);
            font-size: .58em;
            max-height: 52vh;
            overflow: auto;
        }}
        .reveal pre code {{
            white-space: pre-wrap !important;
            word-break: break-word;
        }}
        .reveal code {{
            font-family: 'Consolas', 'JetBrains Mono', monospace;
        }}
    </style>
</head>
<body>
    <div class=\"reveal\">
        <div class=\"slides\">
            <section data-markdown>
                <textarea data-template>
{markdown_content}
                </textarea>
            </section>
        </div>
    </div>
    <script src=\"https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/reveal.js\"></script>
    <script src=\"https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/plugin/markdown/markdown.js\"></script>
    <script src=\"https://cdnjs.cloudflare.com/ajax/libs/reveal.js/4.5.0/plugin/highlight/highlight.js\"></script>
    <script>
        Reveal.initialize({{
            hash: true,
            slideNumber: true,
            center: false,
            controlsLayout: 'edges',
            progress: true,
            transition: 'slide',
            backgroundTransition: 'fade',
            width: 1366,
            height: 768,
            margin: 0.04,
            plugins: [ RevealMarkdown, RevealHighlight ]
        }});
    </script>
</body>
</html>"""
    return html_template.encode("utf-8")


def generate_pptx(md_text: str, topic: str, theme: str = "teaching_blue") -> bytes:
    cleaned_md = sanitize_markdown_for_ppt(md_text)
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    _add_cover_slide(presentation, topic, theme=theme)

    _, sections = _split_markdown_sections(cleaned_md, topic)

    for section in sections:
        section_title = str(section.get("title") or "讲义内容")
        lines = section.get("lines") or ["暂无可解析内容"]
        chunks = _chunk_lines_for_slides(lines, max_weight=12.6, max_code_lines=14)

        for chunk_index, chunk_lines in enumerate(chunks):
            _add_content_slide(
                presentation,
                section_title,
                chunk_lines,
                continued=chunk_index > 0,
                theme=theme,
            )

    bio = io.BytesIO()
    presentation.save(bio)
    return bio.getvalue()


def render_homework_engine():
    render_homework_hero()

    if "homework_engine_last_output" not in st.session_state:
        st.session_state.homework_engine_last_output = ""
    if "homework_engine_last_source" not in st.session_state:
        st.session_state.homework_engine_last_source = ""

    with st.container():
        st.markdown("<div class='homework-shell-hook'></div>", unsafe_allow_html=True)
        st.markdown("<div class='preset-title'>图像输入与推理任务</div>", unsafe_allow_html=True)
        col1, col2 = st.columns([1.15, 2.45], gap="large")
        with col1:
            st.markdown("#### 图像采集")
            upload_img = st.file_uploader(
                "上传题目/作答图片",
                type=["png", "jpg", "jpeg", "webp"],
                key="homework_engine_upload",
            )
            st.caption("或使用摄像头直接拍摄")
            cam_img = st.camera_input("摄像头拍照", key="homework_engine_camera")
            target_image = upload_img or cam_img
            if target_image:
                st.image(target_image, caption="待解析图片", use_container_width=True)

        with col2:
            st.markdown("#### 视觉推理任务")
            tab_extract, tab_variant, tab_grade = st.tabs(
                ["原题数字化", "举一反三", "作答批改"]
            )

            if not target_image:
                st.info("请先在左侧上传或拍摄图片。")
                return

            img_bytes = target_image.getvalue()
            img_suffix = _guess_uploaded_image_suffix(target_image)

            with tab_extract:
                st.markdown("将图中题目转成结构化 Markdown，并给出步骤解析。")
                if st.button("开始提取与解析", use_container_width=True, key="homework_extract_btn"):
                    prompt = (
                        "你是特级教研教师。请读取图片中的题干并输出：\n"
                        "## 1. 数字化原题（公式用 LaTeX）\n"
                        "## 2. 分步解题过程\n"
                        "## 3. 最终答案与易错点提醒"
                    )
                    try:
                        with st.spinner("正在执行视觉识别与推理..."):
                            result = process_image_with_vision_prompt(img_bytes, prompt, suffix=img_suffix)
                        st.session_state.homework_engine_last_output = result
                        st.session_state.homework_engine_last_source = "原题数字化"
                    except Exception as exc:
                        st.error(f"解析失败：{exc}")

            with tab_variant:
                difficulty = st.select_slider(
                    "变式难度",
                    options=["基础巩固", "对标原题", "名校选拔", "奥赛冲刺"],
                    value="对标原题",
                    key="homework_variant_difficulty",
                )
                if st.button("生成同源变式题组", use_container_width=True, key="homework_variant_btn"):
                    prompt = (
                        f"你是命题老师。请识别图中题目考点，生成 3 道【{difficulty}】难度变式题。\n"
                        "输出必须包含：\n"
                        "1) 原题考点拆解\n2) 每道变式题题干\n3) 每道题详细解析\n4) 标准答案"
                    )
                    try:
                        with st.spinner("正在生成变式题组..."):
                            result = process_image_with_vision_prompt(img_bytes, prompt, suffix=img_suffix)
                        st.session_state.homework_engine_last_output = result
                        st.session_state.homework_engine_last_source = f"变式生成-{difficulty}"
                    except Exception as exc:
                        st.error(f"生成失败：{exc}")

            with tab_grade:
                st.markdown("沿学生作答过程定位错误步骤，并给出引导式反馈。")
                if st.button("启动智能批改", use_container_width=True, key="homework_grade_btn"):
                    prompt = (
                        "你是耐心的阅卷教师。请根据图片中的学生作答，输出：\n"
                        "1) 解题思路复盘\n"
                        "2) 逐步核对（标出正确/错误步骤）\n"
                        "3) 错误原因与修正建议\n"
                        "4) 最终评语（鼓励式）"
                    )
                    try:
                        with st.spinner("正在批改并生成反馈..."):
                            result = process_image_with_vision_prompt(img_bytes, prompt, suffix=img_suffix)
                        st.session_state.homework_engine_last_output = result
                        st.session_state.homework_engine_last_source = "作答批改"
                    except Exception as exc:
                        st.error(f"批改失败：{exc}")

    last_output = (st.session_state.homework_engine_last_output or "").strip()
    if last_output:
        with st.container():
            st.markdown("<div class='result-shell-hook'></div>", unsafe_allow_html=True)
            render_result_header(st.session_state.homework_engine_last_source or "智能作答与批改")
            st.markdown(last_output)

            source_name = st.session_state.homework_engine_last_source or "智能作答与批改引擎"
            if st.button("一键推入个人知识库（长期）", use_container_width=True, key="homework_push_to_kb_btn"):
                added_chunks = add_text_to_teacher_kb(last_output, source_name=f"多模态-{source_name}")
                if added_chunks > 0:
                    st.success(f"已入库，新增 {added_chunks} 个知识块。后续讲义生成可检索到这次结果。")
                else:
                    st.warning("当前输出为空，未写入知识库。")
# ==================== 主入口程序 ====================

st.set_page_config(page_title="深学讲义 Agent", layout="wide", initial_sidebar_state="collapsed")
inject_app_styles()

if "app_mode" not in st.session_state:
    st.session_state.app_mode = "讲义生成引擎（备课）"
mode_options = ["讲义生成引擎（备课）", "智能作答与批改（答疑）"]
if st.session_state.app_mode not in mode_options:
    st.session_state.app_mode = "讲义生成引擎（备课）"
if "app_mode_selector" not in st.session_state:
    st.session_state.app_mode_selector = st.session_state.app_mode
if st.session_state.app_mode_selector not in mode_options:
    st.session_state.app_mode_selector = st.session_state.app_mode

with st.container():
    st.markdown("<div class='mode-shell-hook'></div>", unsafe_allow_html=True)
    st.markdown("<div class='preset-title'>核心引擎切换</div>", unsafe_allow_html=True)
    st.session_state.app_mode = st.radio(
        "模式切换",
        mode_options,
        horizontal=True,
        key="app_mode_selector",
        label_visibility="collapsed",
    )
    st.caption("答疑模式可将图像解析结果沉淀到个人知识库。")

if st.session_state.app_mode == "智能作答与批改（答疑）":
    render_homework_engine()
    st.stop()

if "topic_input" not in st.session_state:
    st.session_state.topic_input = ""
if "lesson_result" not in st.session_state:
    st.session_state.lesson_result = None
if "include_code" not in st.session_state:
    st.session_state.include_code = True
if "include_quiz" not in st.session_state:
    st.session_state.include_quiz = True
if "selected_ppt_theme" not in st.session_state:
    st.session_state.selected_ppt_theme = "teaching_blue"
if "teaching_traits" not in st.session_state:
    st.session_state.teaching_traits = ""
if "last_temp_parse_statuses" not in st.session_state:
    st.session_state.last_temp_parse_statuses = []
if "last_kb_parse_statuses" not in st.session_state:
    st.session_state.last_kb_parse_statuses = []

render_hero()

preset_topics = [
    "最新 Python 3.13 新特性与迁移建议",
    "Transformer 注意力机制与可视化教学",
    "FastAPI + Pydantic 实战接口设计",
    "Pandas 数据清洗常见坑与优化策略",
]

with st.container():
    st.markdown("<div class='preset-shell-hook'></div>", unsafe_allow_html=True)
    st.markdown("<div class='preset-title'>快速主题模板（点击即可填入输入框）</div>", unsafe_allow_html=True)
    preset_cols = st.columns(4, gap="small")
    for idx, preset in enumerate(preset_topics):
        with preset_cols[idx]:
            if st.button(preset, key=f"preset_topic_{idx}", use_container_width=True):
                st.session_state.topic_input = preset

uploaded_files = []
persistent_context = ""

with st.container():
    st.markdown("<div class='prompt-shell-hook'></div>", unsafe_allow_html=True)
    with st.expander("高级配置：讲义形态 + 私有知识库", expanded=True):
        st.session_state.teaching_traits = st.text_area(
            "教学风格与硬性要求（可选）",
            value=st.session_state.teaching_traits,
            height=80,
            placeholder="例如：偏案例驱动、面向高一学生、强调板书节奏、每节课都要有提问互动。",
        )

        opt_col1, opt_col2 = st.columns(2, gap="small")
        with opt_col1:
            st.session_state.include_code = st.checkbox(
                "包含示例代码（Python）",
                value=st.session_state.include_code,
                help="关闭后将不生成示例代码章节，也不会执行代码审计。",
            )
        with opt_col2:
            st.session_state.include_quiz = st.checkbox(
                "包含课后选择题与答案",
                value=st.session_state.include_quiz,
                help="关闭后将不生成选择题与答案章节。",
            )

        tabs = st.tabs(["本次临时资料", "个人知识库（长期）"])
        with tabs[0]:
            uploaded_files = st.file_uploader(
                "上传本次要参考的资料",
                type=["txt", "md", "pdf", "docx", "pptx", "html", "htm", "png", "jpg", "jpeg", "webp"],
                accept_multiple_files=True,
            )
            if uploaded_files:
                st.caption(f"本次临时资料 {len(uploaded_files)} 份：{', '.join(f.name for f in uploaded_files)}")
                if st.button("预检临时资料解析状态", use_container_width=True, key="temp_parse_probe_btn"):
                    with st.spinner("正在预检临时资料解析..."):
                        _, _, parse_statuses = build_external_knowledge_context(uploaded_files)
                    st.session_state.last_temp_parse_statuses = parse_statuses
                render_parse_status_panel(st.session_state.get("last_temp_parse_statuses", []), "最近一次临时资料解析状态")

        with tabs[1]:
            current_chunks = teacher_kb_chunk_count()
            kb_info_col, kb_clear_col = st.columns([3.4, 1.2], gap="small")
            with kb_info_col:
                st.caption(f"当前已存知识块：{current_chunks}")
            with kb_clear_col:
                if st.button("清空个人知识库", use_container_width=True, key="teacher_kb_clear_btn"):
                    removed_count, clear_message = clear_teacher_kb()
                    st.session_state.last_kb_parse_statuses = []
                    if "失败" in clear_message:
                        st.error(clear_message)
                    else:
                        st.success(f"{clear_message} 本次共清理约 {removed_count} 个知识块。")
            kb_files = st.file_uploader(
                "添加到个人知识库（可多次追加）",
                type=["txt", "md", "pdf", "docx", "pptx", "html", "htm", "png", "jpg", "jpeg", "webp"],
                accept_multiple_files=True,
                key="teacher_kb_uploader",
            )
            if st.button("入库并建立检索", use_container_width=True):
                if not kb_files:
                    st.warning("请先上传要入库的资料。")
                else:
                    with st.spinner("正在切分并入库资料..."):
                        file_count, chunk_count, parse_statuses = add_files_to_teacher_kb(kb_files)
                    st.session_state.last_kb_parse_statuses = parse_statuses
                    if file_count > 0:
                        st.success(f"入库完成：{file_count} 份资料，新增 {chunk_count} 个知识块。")
                    else:
                        st.warning("资料内容为空或解析失败，未写入知识库。")
            render_parse_status_panel(st.session_state.get("last_kb_parse_statuses", []), "最近一次入库解析状态")

            st.markdown("##### 检索自测")
            kb_probe_query = st.text_input(
                "输入一个关键词，验证长期知识库是否能召回",
                key="kb_probe_query",
                placeholder="例如：红楼梦 最后一回 宝玉出家",
            )
            if st.button("测试知识库召回", use_container_width=True, key="kb_probe_btn"):
                probe = (kb_probe_query or "").strip()
                if not probe:
                    st.warning("请先输入检索关键词。")
                else:
                    hits = query_teacher_kb_hits(probe, n_results=4)
                    if not hits:
                        st.info("当前未命中内容。请确认已入库，或更换关键词后重试。")
                    else:
                        st.success(f"命中 {len(hits)} 条。以下为召回预览：")
                        for i, hit in enumerate(hits, start=1):
                            src = hit.get("source", "未知来源")
                            txt = (hit.get("text", "") or "").strip()
                            score = hit.get("score")
                            score_text = f"（关键词匹配分：{score}）" if isinstance(score, int) else ""
                            preview = _extract_preview_around_query(txt, probe, width=220)
                            st.markdown(f"{i}. **{src}** {score_text}\n\n{preview}...")

    topic_col, action_col = st.columns([5.3, 1.7], gap="small")
    with topic_col:
        st.text_input(
            "topic_input",
            key="topic_input",
            label_visibility="collapsed",
            placeholder="输入教学核心主题，例如：最新 Python 更新、RAG 架构设计、机器学习模型评估",
        )
    with action_col:
        generate_clicked = st.button("生成发行级讲义", use_container_width=True)

if generate_clicked:
    topic = st.session_state.topic_input.strip()
    if not topic:
        st.warning("请输入教学主题后再生成。")
    else:
        temp_knowledge, knowledge_sources, temp_parse_statuses = build_external_knowledge_context(
            uploaded_files or [],
            query_text=topic,
        )
        st.session_state.last_temp_parse_statuses = temp_parse_statuses
        kb_hits = query_teacher_kb_hits(topic, n_results=4)
        persistent_context = "\n\n".join(
            f"[知识库:{hit.get('source', '未知来源')}] {(hit.get('text', '') or '').strip()}"
            for hit in kb_hits
            if (hit.get("text") or "").strip()
        )
        external_knowledge = "\n\n".join(x for x in [temp_knowledge, persistent_context] if x.strip())
        if kb_hits:
            knowledge_sources = [*knowledge_sources, f"个人知识库检索结果（命中{len(kb_hits)}条）"]

        status_box = st.status("正在调度多智能体协同处理...", expanded=True)
        initial_state = {
            "topic": topic,
            "iteration_count": 0,
            "messages": [],
            "teaching_traits": st.session_state.teaching_traits,
            "external_knowledge": external_knowledge,
            "knowledge_sources": knowledge_sources,
            "include_code": st.session_state.include_code,
            "include_quiz": st.session_state.include_quiz,
            "repair_target": "",
            "feedback_category": "",
        }
        final_state = initial_state
        review_history: list[dict] = []
        repair_labels = {
            "subject": "教研 Agent",
            "pedagogy": "教学设计 Agent",
            "activity": "互动设计 Agent",
            "quiz": "题库 Agent",
            "content": "总编 Agent",
        }

        try:
            for event in app_workflow.stream(initial_state, config={"recursion_limit": 50}):
                for node_name, state_update in event.items():
                    with status_box:
                        if node_name == "SubjectExpert" and knowledge_sources:
                            st.write(f"外部知识库：已注入 {len(knowledge_sources)} 份教师资料。")
                        if node_name == "SubjectExpert":
                            if final_state.get("repair_target") == "subject":
                                st.write("教研 Agent：收到质检反馈，正在重做主题范围与学科大纲。")
                            else:
                                st.write("教研 Agent：执行主题拆解、知识检索与大纲构建。")
                        elif node_name == "PedagogyPlanner":
                            if final_state.get("repair_target") == "pedagogy":
                                st.write("教学设计 Agent：收到质检反馈，正在重做课堂流程、学情假设与分层策略。")
                            else:
                                st.write("教学设计 Agent：生成课堂流程、学情假设与分层策略。")
                        elif node_name == "ActivityDesigner":
                            if final_state.get("repair_target") == "activity":
                                st.write("互动设计 Agent：收到质检反馈，正在重做课堂活动与过程性评价。")
                            else:
                                st.write("互动设计 Agent：编排课堂活动与过程性评价。")
                        elif node_name == "QuizDesigner":
                            if final_state.get("include_quiz", True):
                                if final_state.get("repair_target") == "quiz":
                                    st.write("题库 Agent：收到质检反馈，正在重做完整题目、选项、答案与解析。")
                                else:
                                    st.write("题库 Agent：生成分层选择题与答案解析。")
                            else:
                                st.write("题库 Agent：本次已关闭，跳过题库生成。")
                        elif node_name == "ContentCreator":
                            current_iter = state_update.get("iteration_count", 0)
                            if final_state.get("repair_target") and final_state.get("repair_target") != "content":
                                st.write(f"总编 Agent：整合已返工模块并输出完整讲义（第 {current_iter} 轮）。")
                            elif final_state.get("repair_target") == "content":
                                st.write(f"总编 Agent：根据质检反馈重写完整讲义（第 {current_iter} 轮）。")
                            else:
                                st.write(f"总编 Agent：融合多源结果并输出讲义（第 {current_iter} 轮）。")
                        elif node_name == "EvaluationExpert":
                            if state_update.get("is_approved"):
                                st.write("质检 Agent：模板、题库和代码检查全部通过。")
                            else:
                                target = state_update.get("repair_target") or "content"
                                category = state_update.get("feedback_category") or "质量问题"
                                target_label = repair_labels.get(target, "总编 Agent")
                                feedback = (state_update.get("feedback") or "").strip()
                                debug_info = _collect_runtime_review_debug(
                                    final_state.get("draft_content", ""),
                                    include_quiz=final_state.get("include_quiz", True),
                                )
                                review_history.append(
                                    {
                                        "round": final_state.get("iteration_count", 0),
                                        "target": target,
                                        "category": category,
                                        "feedback": feedback,
                                        "debug": debug_info,
                                    }
                                )
                                st.write(f"质检 Agent：发现{category}，已路由给{target_label}返工。")
                                if feedback:
                                    st.caption("质检原文：" + feedback)
                                if target == "quiz":
                                    st.caption(
                                        f"题库诊断：题目 {debug_info.get('question_count', 0)} 条，"
                                        f"答案解析 {debug_info.get('answer_count', 0)} 条。"
                                    )
                    final_state.update(state_update)

            include_code = final_state.get("include_code", st.session_state.include_code)
            include_quiz = final_state.get("include_quiz", st.session_state.include_quiz)
            raw_md = (final_state.get("final_content") or "").strip()
            if not final_state.get("is_approved") or not raw_md:
                feedback = (final_state.get("feedback") or "质检未给出具体原因。").strip()
                last_draft = (final_state.get("draft_content") or "").strip()
                if not last_draft:
                    status_box.update(label="讲义未通过质检", state="error", expanded=True)
                    st.session_state.lesson_result = None
                    st.error("本次生成未通过质检，且无可用草稿。请重试。")
                    with st.expander("查看最后一次质检反馈", expanded=True):
                        st.write(feedback)
                    st.stop()

                status_box.update(label="讲义已生成完成（自动修复兜底）", state="complete", expanded=False)
                st.warning("本次未完全通过质检，系统已自动采用最后一版草稿并执行模板修复后继续交付。")
                with st.expander("查看最后一次质检反馈", expanded=False):
                    st.write(feedback)
                raw_md = last_draft
            else:
                status_box.update(label="讲义已生成完成", state="complete", expanded=False)

            final_md = ensure_markdown_template(raw_md, topic, include_code=include_code, include_quiz=include_quiz)
            quality_issues = _lesson_output_quality_issues(final_md, include_code=include_code, include_quiz=include_quiz)
            recovered_by_local_builder = False
            if quality_issues:
                with status_box:
                    st.write("本地质量闸门：模型最终草稿不完整，正在用各子 Agent 结果重组讲义。")
                recovery_raw = _build_recovery_lesson_markdown(topic, final_state, include_code, include_quiz)
                final_md = ensure_markdown_template(recovery_raw, topic, include_code=include_code, include_quiz=include_quiz)
                recovered_by_local_builder = True
                quality_issues = _lesson_output_quality_issues(final_md, include_code=include_code, include_quiz=include_quiz)
            code_audit = prepare_code_runner_agent(final_md) if include_code else None

            st.session_state.lesson_result = {
                "topic": topic,
                "final_md": final_md,
                "generated_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                "code_audit": code_audit,
                "knowledge_sources": knowledge_sources,
                "kb_hits_preview": kb_hits,
                "include_code": include_code,
                "include_quiz": include_quiz,
                "teaching_traits": st.session_state.teaching_traits,
                "recovered_by_local_builder": recovered_by_local_builder,
                "quality_issues": quality_issues,
                "review_history": review_history,
            }

        except Exception as e:
            status_box.update(label="执行流发生中断", state="error")
            st.error(f"引擎异常信息: {str(e)}")

result_payload = st.session_state.lesson_result

if result_payload:
    topic = result_payload["topic"]
    include_code = result_payload.get("include_code", st.session_state.include_code)
    include_quiz = result_payload.get("include_quiz", st.session_state.include_quiz)
    final_md = ensure_markdown_template(
        result_payload["final_md"],
        topic,
        include_code=include_code,
        include_quiz=include_quiz,
    )
    st.session_state.lesson_result["final_md"] = final_md
    knowledge_sources = result_payload.get("knowledge_sources") or []
    
    code_audit = result_payload.get("code_audit")
    if include_code and (not code_audit or "status" not in code_audit):
        code_audit = prepare_code_runner_agent(final_md)
        st.session_state.lesson_result["code_audit"] = code_audit

    with st.container():
        st.markdown("<div class='result-shell-hook'></div>", unsafe_allow_html=True)
        render_result_header(topic)
        traits_used = (result_payload.get("teaching_traits") or "").strip()
        if traits_used:
            st.caption(f"本次教学风格约束：{traits_used}")
        if result_payload.get("recovered_by_local_builder"):
            st.warning("模型最终草稿不完整，系统已使用大纲、教学设计、互动活动和题库草案重组讲义。")
        residual_issues = result_payload.get("quality_issues") or []
        if residual_issues:
            st.error("当前讲义仍存在质量问题：" + "；".join(residual_issues))
        review_history = result_payload.get("review_history") or []
        if review_history:
            with st.expander("查看完整质检历史", expanded=False):
                for item in review_history:
                    round_no = item.get("round", "?")
                    category = item.get("category", "质量问题")
                    target = item.get("target", "content")
                    feedback = (item.get("feedback") or "无反馈原文").strip()
                    debug = item.get("debug") or {}
                    st.markdown(f"**第 {round_no} 轮** · {category} · 路由到 `{target}`")
                    st.write(feedback)
                    if debug:
                        st.caption(
                            f"识别结果：题目 {debug.get('question_count', 0)} 条，"
                            f"答案解析 {debug.get('answer_count', 0)} 条"
                        )
                        if debug.get("quiz_preview"):
                            st.code(debug.get("quiz_preview"), language="markdown")
                        if debug.get("answer_preview"):
                            st.code(debug.get("answer_preview"), language="markdown")
        if knowledge_sources:
            st.caption(f"本次生成已参考外部资料：{', '.join(knowledge_sources)}")
        kb_hits_preview = result_payload.get("kb_hits_preview") or []
        if kb_hits_preview:
            st.caption(f"个人知识库命中 {len(kb_hits_preview)} 条并已注入生成上下文。")
            with st.expander("查看个人知识库命中片段（前 4 条）", expanded=False):
                for i, hit in enumerate(kb_hits_preview[:4], start=1):
                    src = hit.get("source", "未知来源")
                    txt = (hit.get("text", "") or "").strip()
                    preview = _extract_preview_around_query(txt, topic, width=260)
                    st.markdown(f"{i}. **{src}**\n\n{preview}...")
        st.markdown(final_md)

    if include_code:
        with st.container():
            st.markdown("<div class='code-audit-hook'></div>", unsafe_allow_html=True)
            is_pending = code_audit.get("status") == "pending"

            if is_pending:
                summary_text = f"blocks={code_audit.get('total', 0)} | 等待验证状态"
                st.info(f"**运行环境确认**: 本地使用的 Python 环境路径为 `{sys.executable}`，执行前请确保可能需要的第三方库已经安装。")
            else:
                summary_text = (
                    f"blocks={code_audit.get('total', 0)} | "
                    f"pass={code_audit.get('passed', 0)} | "
                    f"fail={code_audit.get('failed', 0)} | "
                    f"time={code_audit.get('total_duration_ms', 0)}ms"
                )

            st.markdown(
                f"""
                <div class="code-audit-head">
                    <div class="code-audit-title">代码运行 Agent 报告</div>
                    <div class="code-audit-meta">{summary_text}</div>
                </div>
                """,
                unsafe_allow_html=True,
            )

            tool_col, download_col = st.columns([1.1, 1], gap="small")
            with tool_col:
                btn_label = "开始执行验证 Code Runner Agent" if is_pending else "重新执行代码块"
                if st.button(btn_label, key="rerun_code_agent", use_container_width=True):
                    with st.spinner("Code Runner Agent 正在执行内部审计与依赖检查..."):
                        updated_audit = run_code_runner_agent(final_md)
                        st.session_state.lesson_result["code_audit"] = updated_audit
                        code_audit = updated_audit
                    st.success("代码运行完毕，审计报告已更新。您可以展开下方手风琴查看详细测试与安装输出。")
                    st.rerun()

            with download_col:
                audit_md = build_code_audit_markdown(topic, code_audit)
                st.download_button(
                    "下载代码运行报告",
                    data=audit_md,
                    file_name=f"{topic}_代码运行报告.md",
                    mime="text/markdown",
                    use_container_width=True,
                )

            if code_audit.get("total", 0) == 0:
                st.info("当前讲义未检测到 python 代码块。")
            else:
                for item in code_audit.get("reports", []):
                    if item.get("is_pending"):
                        status = "PENDING"
                        expanded = True
                    else:
                        status = "PASS" if item.get("success") else "FAIL"
                        expanded = not item.get("success")
                    block_label = (
                        f"代码块 {item.get('index')} · {status} · {item.get('duration_ms', 0)}ms"
                    )
                    with st.expander(block_label, expanded=expanded):
                        st.caption("分析的代码内容：请注意前置引用的包")
                        st.code(item.get("code", ""), language="python")

                        if item.get("is_pending"):
                            st.caption("安装提示 (供您自行验证参考，该 Agent 不会自发变更您的 Python 全局环境，除非您点击验证按钮)")
                            imported_modules = item.get("imported_modules", [])
                            if imported_modules:
                                pip_cmd = f"{sys.executable} -m pip install " + " ".join(imported_modules)
                                hint_text = f"# 引擎检测到代码包含外部模块，若在本地验证请先执行：\n{pip_cmd}"
                            else:
                                hint_text = "# 未检测到明确的外部依赖。如果运行时仍然缺少依赖，再使用 pip install 进行安装。"

                            st.code(hint_text, language="bash")
                        else:
                            st.caption("执行输出 / 安装日志")
                            st.code(item.get("output", ""), language="text")
    else:
        st.info("本次已关闭示例代码生成功能，因此跳过代码运行审计。")

    st.markdown("<div class='export-title'>Export Bundle</div>", unsafe_allow_html=True)
    if st.session_state.selected_ppt_theme not in PPT_THEMES:
        st.session_state.selected_ppt_theme = "teaching_blue"
    if "ppt_theme_picker" not in st.session_state or st.session_state.ppt_theme_picker not in PPT_THEMES:
        st.session_state.ppt_theme_picker = st.session_state.selected_ppt_theme

    theme_keys = list(PPT_THEMES.keys())
    selected_theme_key = st.radio(
        "PPT 视觉主题",
        options=theme_keys,
        format_func=lambda key: PPT_THEMES[key]["name"],
        horizontal=True,
        key="ppt_theme_picker",
    )
    st.session_state.selected_ppt_theme = selected_theme_key
    selected_theme_label = PPT_THEMES[selected_theme_key]["name"]

    col1, col2, col3, col4, col5 = st.columns(5, gap="small")
    with col1:
        st.download_button(
            "下载 Markdown",
            data=final_md,
            file_name=f"{topic}_讲义.md",
            mime="text/markdown",
        )
    with col2:
        docx_bytes = generate_docx(final_md)
        st.download_button(
            "下载 Word",
            data=docx_bytes,
            file_name=f"{topic}_讲义.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )
    with col3:
        pdf_bytes = generate_pdf(final_md)
        if pdf_bytes is not None:
            st.download_button(
                "下载 PDF",
                data=pdf_bytes,
                file_name=f"{topic}_讲义.pdf",
                mime="application/pdf",
            )
        else:
            st.button("PDF 环境未加载", disabled=True, key="pdf_disabled")
    with col4:
        ppt_ready_md = sanitize_markdown_for_ppt(final_md)
        pptx_bytes = generate_pptx(
            ppt_ready_md,
            topic,
            theme=st.session_state.selected_ppt_theme,
        )
        st.download_button(
            "下载 PPT",
            data=pptx_bytes,
            file_name=f"{topic}_讲义_{selected_theme_label}.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
    with col5:
        html_bytes = generate_html_presentation(ppt_ready_md, topic)
        st.download_button(
            "下载 HTML",
            data=html_bytes,
            file_name=f"{topic}_讲义演示.html",
            mime="text/html",
        )
else:
    st.info("输入教学主题后点击“生成发行级讲义”，系统将自动完成大纲、讲义、代码和题库质检。")

